#!/usr/bin/env python3
# -*- coding: utf-8 -*-
# NinjaEssayAI Bot - Automated Essay Writing Service

import os
import asyncio
import io
import sqlite3
import time
import csv
import sys
import logging
import aiohttp
from collections import defaultdict

# Настройка кодировки для Windows консоли
if sys.platform == 'win32':
    import codecs
    sys.stdout = codecs.getwriter('utf-8')(sys.stdout.buffer, 'strict')
    sys.stderr = codecs.getwriter('utf-8')(sys.stderr.buffer, 'strict')

try:
    import psutil
except ImportError:
    psutil = None
from telegram import Update, ReplyKeyboardMarkup, InlineKeyboardButton, InlineKeyboardMarkup
from telegram.ext import ApplicationBuilder, CommandHandler, CallbackContext, MessageHandler, filters, ConversationHandler, CallbackQueryHandler
from openai import AsyncOpenAI
import docx
from docx.shared import Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
from docx.enum.section import WD_SECTION
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
import logging
import json
import re
import html
from dotenv import load_dotenv
from sqlalchemy import create_engine, Column, Integer, String, DateTime
from sqlalchemy.orm import sessionmaker, declarative_base  # Updated import for SQLAlchemy 2.0
from datetime import datetime, timezone, timedelta
import uuid
from yookassa import Configuration, Payment, Refund

# Импорты для презентаций
from pptx import Presentation
from pptx.util import Inches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from PIL import Image
import requests
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')  # Для работы без GUI

# Загрузка переменных окружения
load_dotenv()
TELEGRAM_BOT_TOKEN = os.getenv("TELEGRAM_BOT_TOKEN")
DEEPSEEK_API_KEY = os.getenv("DEEPSEEK_API_KEY")
DEEPSEEK_BASE_URL = os.getenv("DEEPSEEK_BASE_URL", "https://api.deepseek.com")
# Значения по умолчанию для YooKassa удалены для безопасности
YOOKASSA_SHOP_ID = os.getenv("YOOKASSA_SHOP_ID")
YOOKASSA_SECRET_KEY = os.getenv("YOOKASSA_SECRET_KEY")

# Настройки Coze API для поиска источников
COZE_API_TOKEN = os.getenv("COZE_API_TOKEN")
COZE_WORKFLOW_ID = os.getenv("COZE_WORKFLOW_ID")
COZE_SPACE_ID = os.getenv("COZE_SPACE_ID")
COZE_API_URL = os.getenv("COZE_API_URL", "https://api.coze.com/v1/workflow/run")

# 🧪 ТЕСТОВЫЙ РЕЖИМ - установите в False для рабочего режима с реальными платежами
TESTING_MODE = False

# Администраторы бота
ADMIN_IDS = [659874549]  # Ваш Telegram ID как администратор

# Пользователи с постоянным тестовым режимом (по username)
# Для этих пользователей оплата всегда бесплатная
TEST_MODE_USERNAMES = ["AbuHavva"]  # Username без @

# Система ограничения запросов (rate limiting)
user_request_times = defaultdict(list)
MAX_REQUESTS_PER_HOUR = 5
REQUEST_WINDOW = 3600  # 1 час в секундах

# Настройка YooKassa
Configuration.account_id = YOOKASSA_SHOP_ID or ""
Configuration.secret_key = YOOKASSA_SECRET_KEY or ""

# Проверяем, что ключ DeepSeek загружен корректно
if not DEEPSEEK_API_KEY:
    raise ValueError(
        "Переменная окружения DEEPSEEK_API_KEY не установлена. "
        "Убедитесь, что файл .env содержит корректный ключ API."
    )
# Предупреждаем, если переменные YooKassa не заданы
if not YOOKASSA_SHOP_ID or not YOOKASSA_SECRET_KEY:
    logging.warning(
        "Переменные окружения YOOKASSA_SHOP_ID или YOOKASSA_SECRET_KEY "
        "не установлены. Платежи через YooKassa не будут работать."
    )
# Предупреждаем, если переменные Coze API не заданы
if not COZE_API_TOKEN or not COZE_WORKFLOW_ID or not COZE_SPACE_ID:
    logging.warning(
        "Переменные окружения COZE_API_TOKEN, COZE_WORKFLOW_ID или COZE_SPACE_ID "
        "не установлены. Автоматический поиск источников не будет работать."
    )

client = AsyncOpenAI(api_key=DEEPSEEK_API_KEY, base_url=DEEPSEEK_BASE_URL)

# Настройка логирования с правильной кодировкой для Windows
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler(sys.stdout)
    ]
)
# Явно устанавливаем кодировку для логов
for handler in logging.root.handlers:
    if isinstance(handler, logging.StreamHandler):
        handler.stream = sys.stdout

# Функции валидации и безопасности
def sanitize_filename(text: str) -> str:
    """Безопасная очистка имени файла"""
    if not text:
        return "default"
    # Убираем опасные символы
    safe_text = re.sub(r'[<>:"/\\|?*\x00-\x1f]', '_', text)
    # Ограничиваем длину
    safe_text = safe_text[:50].strip()
    # Убираем точки в начале и конце (скрытые файлы/папки)
    safe_text = safe_text.strip('.')
    return safe_text if safe_text else "default"

def validate_user_input(text: str, max_length: int = 1000) -> str:
    """Валидация пользовательского ввода"""
    if not text:
        raise ValueError("Пустой ввод не допускается")
    if len(text) > max_length:
        raise ValueError(f"Слишком длинный текст (максимум {max_length} символов)")
    # Убираем потенциально опасные теги и экранируем HTML
    cleaned_text = html.escape(text.strip())
    if not cleaned_text:  # Проверяем после strip()
        raise ValueError("Пустой ввод не допускается")
    return cleaned_text

def validate_contact(contact: str) -> str:
    """Валидация контактных данных"""
    contact = validate_user_input(contact, 100)
    # Дополнительная валидация для email/телефона
    pattern_email = re.compile(r"^[\w\.-]+@[\w\.-]+\.\w+$")
    # Паттерн для телефона: цифры, пробелы, тире, плюс в начале
    pattern_phone = re.compile(r"^\+?[0-9\s\-\(\)]{10,20}$")
    if not (pattern_email.match(contact) or pattern_phone.match(contact)):
        raise ValueError("Неверный формат контакта")
    return contact

def get_chat_id(context: CallbackContext, update: Update = None) -> int:
    """Безопасное получение chat_id"""
    if hasattr(context, '_chat_id') and context._chat_id:
        return context._chat_id
    elif update and update.effective_chat:
        return update.effective_chat.id
    elif update and update.message:
        return update.message.chat_id
    else:
        raise ValueError("Не удалось определить chat_id")

def remove_emojis(text: str) -> str:
    """Удаляет все эмодзи и смайлики из текста"""
    if not text:
        return text
    
    # Расширенный паттерн для удаления всех эмодзи
    emoji_pattern = re.compile(
        "["
        u"\U0001F600-\U0001F64F"  # эмоции
        u"\U0001F300-\U0001F5FF"  # символы и пиктограммы
        u"\U0001F680-\U0001F6FF"  # транспорт и символы на карте
        u"\U0001F1E0-\U0001F1FF"  # флаги
        u"\U00002702-\U000027B0"  # разное
        u"\U000024C2-\U0001F251"  # заключенные символы
        u"\U0001F900-\U0001F9FF"  # дополнительные смайлики
        u"\U0001FA70-\U0001FAFF"  # расширенные символы
        u"\U00002500-\U00002BEF"  # китайские символы
        u"\U0001F018-\U0001F270"  # разные символы
        u"\U00002300-\U000023FF"  # технические символы
        u"\U0001F004-\U0001F0CF"  # игровые символы
        "]+", flags=re.UNICODE
    )
    return emoji_pattern.sub('', text)

def validate_generated_content(content: str, chapter: str) -> str:
    """Валидирует и очищает сгенерированный контент от нежелательных фраз и смайликов"""
    if not content or len(content.strip()) < 100:
        raise ValueError(f"Слишком короткий контент для главы {chapter}")
    
    # Удаляем все смайлики и эмодзи из текста
    content = remove_emojis(content)
    
    # Удаляем нежелательные фразы в начале текста
    unwanted_patterns = [
        r'^[^.]*отлично[^.]*\.',
        r'^[^.]*вот план[^.]*\.',
        r'^[^.]*рассмотрим[^.]*\.',
        r'^[^.]*вот текст[^.]*\.',
        r'^[^.]*итак[^.]*\.',
        r'^[^.]*составленный с учетом[^.]*\.',
        r'^[^.]*план .* по теме[^.]*\.',
    ]
    
    content_cleaned = content.strip()
    
    # Удаляем первое предложение, если оно содержит нежелательные фразы
    for pattern in unwanted_patterns:
        content_cleaned = re.sub(pattern, '', content_cleaned, flags=re.IGNORECASE)
    
    # Удаляем лишние пробелы и переносы в начале
    content_cleaned = content_cleaned.lstrip()
    
    # Проверяем, что осталось достаточно контента
    if len(content_cleaned.strip()) < 200:
        logging.warning(f"Контент для главы {chapter} стал слишком коротким после очистки")
        return content.strip()  # Возвращаем оригинал
    
    return content_cleaned

def remove_chapter_title_from_text(chapter_text: str, chapter_title: str) -> str:
    """Удаляет дублирующийся заголовок главы из начала текста
    
    Args:
        chapter_text: Текст главы, возможно содержащий заголовок
        chapter_title: Название главы для поиска и удаления
    
    Returns:
        Текст без дублирующегося заголовка
    """
    if not chapter_text or not chapter_title:
        return chapter_text
    
    # Очищаем название главы от номеров и лишних символов для поиска
    clean_title = re.sub(r'^\d+\.\s*', '', chapter_title).strip()
    clean_title = re.sub(r'^Глава\s+\d+[\.:]\s*', '', clean_title, flags=re.IGNORECASE).strip()
    
    # Паттерны для поиска заголовка в начале текста
    patterns = [
        # С номером главы в начале: "1. Введение", "Глава 1. Введение", "Глава 1: Введение"
        rf'^Глава\s+\d+[\.:]\s*{re.escape(clean_title)}\.?\s*\n*',
        rf'^\d+\.\s*{re.escape(clean_title)}\.?\s*\n*',
        # Без номера: "Введение", "ВВЕДЕНИЕ"
        rf'^{re.escape(clean_title)}\.?\s*\n*',
        rf'^{re.escape(clean_title.upper())}\.?\s*\n*',
        # С дополнительными символами
        rf'^#+\s*{re.escape(clean_title)}\.?\s*\n*',
        rf'^\*\*{re.escape(clean_title)}\*\*\.?\s*\n*',
    ]
    
    # Пробуем удалить заголовок по каждому паттерну
    result = chapter_text.strip()
    for pattern in patterns:
        result = re.sub(pattern, '', result, count=1, flags=re.IGNORECASE | re.MULTILINE)
    
    # Удаляем лишние переносы строк в начале
    result = result.lstrip('\n').strip()
    
    return result

# ===================== ФУНКЦИИ ДЛЯ РАБОТЫ С ИСТОЧНИКАМИ =====================

async def fetch_sources_from_coze(keywords: str, count: int = 15) -> list:
    """
    Получает источники через Coze workflow API
    
    Args:
        keywords: Ключевые слова для поиска источников
        count: Количество источников для поиска (не используется, т.к. workflow возвращает фиксированное количество)
    
    Returns:
        Список словарей с ключами 'title' и 'url'
    """
    logging.info(f"Запрос источников через Coze workflow по ключевым словам: {keywords}")
    
    headers = {
        "Authorization": f"Bearer {COZE_API_TOKEN}",
        "Content-Type": "application/json"
    }
    
    # ВАЖНО: Простой запрос (только ключевые слова) даёт максимум источников (~10)
    # Сложные инструкции с указанием количества уменьшают результат до 3 источников
    payload = {
        "workflow_id": COZE_WORKFLOW_ID,
        "parameters": {
            "input": keywords  # Простой запрос без дополнительных инструкций
        }
    }
    
    try:
        async with aiohttp.ClientSession() as session:
            async with session.post(COZE_API_URL, headers=headers, json=payload, timeout=60) as response:
                if response.status == 200:
                    data = await response.json()
                    logging.info(f"Ответ от Coze API: {data}")
                    
                    # Извлекаем источники из ответа
                    sources = parse_coze_response(data)
                    logging.info(f"Получено источников: {len(sources)}")
                    return sources
                else:
                    error_text = await response.text()
                    logging.error(f"Ошибка при запросе к Coze API: {response.status}, {error_text}")
                    return []
    except asyncio.TimeoutError:
        logging.error("Таймаут при запросе к Coze API")
        return []
    except Exception as e:
        logging.error(f"Исключение при запросе к Coze API: {e}")
        return []

def parse_coze_response(data: dict) -> list:
    """
    Парсит ответ от Coze workflow и извлекает источники
    
    Args:
        data: JSON ответ от Coze API
    
    Returns:
        Список словарей с ключами 'title' и 'url'
    """
    sources = []
    
    try:
        # Извлекаем данные из ответа
        if "data" in data:
            data_str = data.get("data", "{}")
            
            # Если data - это строка, парсим как JSON
            if isinstance(data_str, str):
                try:
                    data_obj = json.loads(data_str)
                except json.JSONDecodeError:
                    logging.error("Не удалось распарсить JSON из data")
                    return []
            else:
                data_obj = data_str
            
            # Извлекаем массив output
            if "output" in data_obj:
                output_array = data_obj.get("output", [])
                
                if isinstance(output_array, list):
                    for item in output_array:
                        if isinstance(item, dict):
                            # Извлекаем title и link
                            title = item.get("title", "")
                            link = item.get("link", "")
                            
                            if title and link:
                                sources.append({"title": title, "url": link})
                        elif isinstance(item, str):
                            # Если элемент - строка, парсим её
                            parsed = parse_source_string(item)
                            if parsed:
                                sources.append(parsed)
        
        logging.info(f"Распарсено источников: {len(sources)}")
    except Exception as e:
        logging.error(f"Ошибка при парсинге ответа Coze: {e}")
    
    return sources

def parse_sources_from_text(text: str) -> list:
    """
    Парсит источники из текстового ответа
    
    Args:
        text: Текст с источниками
    
    Returns:
        Список словарей с ключами 'title' и 'url'
    """
    sources = []
    
    # Ищем паттерны вида "Заголовок - URL" или "Заголовок: URL"
    patterns = [
        r'(.+?)\s*[-–—]\s*(https?://[^\s]+)',
        r'(.+?):\s*(https?://[^\s]+)',
        r'\d+\.\s*(.+?)\s*[-–—]\s*(https?://[^\s]+)',
        r'\d+\.\s*(.+?):\s*(https?://[^\s]+)',
    ]
    
    for pattern in patterns:
        matches = re.findall(pattern, text, re.MULTILINE)
        for match in matches:
            title = match[0].strip()
            url = match[1].strip()
            sources.append({"title": title, "url": url})
    
    # Если ничего не нашли, попробуем найти просто URL
    if not sources:
        urls = re.findall(r'https?://[^\s]+', text)
        for i, url in enumerate(urls, 1):
            sources.append({"title": f"Источник {i}", "url": url})
    
    return sources

def parse_source_string(text: str) -> dict:
    """
    Парсит одну строку с источником
    
    Args:
        text: Строка с источником
    
    Returns:
        Словарь с ключами 'title' и 'url' или None
    """
    patterns = [
        r'(.+?)\s*[-–—]\s*(https?://[^\s]+)',
        r'(.+?):\s*(https?://[^\s]+)',
    ]
    
    for pattern in patterns:
        match = re.match(pattern, text.strip())
        if match:
            return {"title": match.group(1).strip(), "url": match.group(2).strip()}
    
    # Если это просто URL
    if text.strip().startswith('http'):
        return {"title": "Источник", "url": text.strip()}
    
    return None

def format_source_gost(source: dict, index: int) -> str:
    """
    Оформляет источник по ГОСТу
    
    Args:
        source: Словарь с ключами 'title' и 'url'
        index: Номер источника в списке
    
    Returns:
        Строка с оформленным источником
    """
    title = source.get("title", "Без названия")
    url = source.get("url", "")
    
    # Базовое оформление по ГОСТ 7.0.5-2008
    # Для электронных ресурсов
    if url:
        # Определяем тип источника по URL
        if "doi.org" in url or "scholar.google" in url or "elibrary.ru" in url:
            # Научная статья
            formatted = f"{index}. {title} [Электронный ресурс]. – URL: {url} (дата обращения: {datetime.now().strftime('%d.%m.%Y')})."
        elif "wikipedia.org" in url:
            # Википедия
            formatted = f"{index}. {title} // Википедия [Электронный ресурс]. – URL: {url} (дата обращения: {datetime.now().strftime('%d.%m.%Y')})."
        else:
            # Обычный веб-ресурс
            formatted = f"{index}. {title} [Электронный ресурс]. – URL: {url} (дата обращения: {datetime.now().strftime('%d.%m.%Y')})."
    else:
        # Если URL отсутствует
        formatted = f"{index}. {title}."
    
    return formatted

def extract_keywords_from_theme(theme: str, science_name: str) -> str:
    """
    Извлекает ключевые слова из темы работы для поиска источников
    
    Args:
        theme: Тема работы
        science_name: Название предмета
    
    Returns:
        Строка с ключевыми словами
    """
    # Объединяем тему и предмет для более точного поиска
    keywords = f"{theme} {science_name}"
    
    # Убираем лишние слова
    stop_words = ["по", "для", "в", "на", "с", "о", "и", "или", "а", "но"]
    words = keywords.split()
    filtered_words = [w for w in words if w.lower() not in stop_words]
    
    return " ".join(filtered_words)

# Функции администрирования и безопасности
def is_admin(user_id: int) -> bool:
    """Проверка, является ли пользователь администратором"""
    return user_id in ADMIN_IDS

def check_rate_limit(user_id: int) -> bool:
    """Проверка лимита запросов пользователя"""
    current_time = time.time()
    user_times = user_request_times[user_id]
    
    # Удаляем старые запросы за пределами окна
    user_times[:] = [t for t in user_times if current_time - t < REQUEST_WINDOW]
    
    # Проверяем лимит
    if len(user_times) >= MAX_REQUESTS_PER_HOUR:
        return False
    
    # Добавляем текущий запрос
    user_times.append(current_time)
    return True

async def admin_stats(update: Update, context: CallbackContext) -> None:
    """Статистика для администраторов"""
    if not is_admin(update.effective_user.id):
        await update.message.reply_text("❌ Доступ запрещен")
        return
    
    try:
        session = SessionLocal()
        
        # Общая статистика
        total_actions = session.query(UserAction).count()
        unique_users = session.query(UserAction.user_id).distinct().count()
        
        # Статистика по действиям
        start_commands = session.query(UserAction).filter(UserAction.action == "start_command").count()
        order_commands = session.query(UserAction).filter(UserAction.action == "order_command").count()
        
        # Статистика за последние 24 часа
        day_ago = datetime.now(timezone.utc) - timedelta(hours=24)
        recent_actions = session.query(UserAction).filter(UserAction.timestamp > day_ago).count()
        
        # Дополнительные метрики
        week_ago = datetime.now(timezone.utc) - timedelta(days=7)
        week_actions = session.query(UserAction).filter(UserAction.timestamp > week_ago).count()
        
        # Топ-3 популярных действий
        from sqlalchemy import func
        top_actions = session.query(
            UserAction.action,
            func.count(UserAction.id).label('count')
        ).group_by(UserAction.action).order_by(func.count(UserAction.id).desc()).limit(3).all()
        
        # Активность по дням недели
        from sqlalchemy import extract
        today_weekday = datetime.now().weekday()
        today_actions = session.query(UserAction).filter(
            extract('dow', UserAction.timestamp) == today_weekday,
            UserAction.timestamp > day_ago
        ).count()
        
        # Новые пользователи за неделю
        new_users_week = session.query(UserAction.user_id).filter(
            UserAction.timestamp > week_ago
        ).distinct().count()
        
        # Активные пользователи (действия за последние 7 дней)
        active_users_week = session.query(UserAction.user_id).filter(
            UserAction.timestamp > week_ago
        ).distinct().count()
        
        # Среднее количество действий на пользователя
        avg_actions_per_user = total_actions / unique_users if unique_users > 0 else 0
        
        # === ДОПОЛНИТЕЛЬНЫЕ МЕТРИКИ ===
        # Статистика заказов
        total_orders = session.query(Order).count()
        paid_orders = session.query(Order).filter(Order.status == "paid").count()
        completed_orders = session.query(Order).filter(Order.status == "completed").count()
        failed_orders = session.query(Order).filter(Order.status == "failed").count()
        
        # Заказы за последние 24 часа и неделю
        recent_orders = session.query(Order).filter(Order.created_at > day_ago).count()
        week_orders = session.query(Order).filter(Order.created_at > week_ago).count()
        
        # Конверсия (% пользователей, сделавших заказ)
        users_with_orders = session.query(Order.user_id).distinct().count()
        conversion_rate = (users_with_orders / unique_users * 100) if unique_users > 0 else 0
        
        # Средняя цена заказа
        from sqlalchemy import func
        avg_order_price = session.query(func.avg(Order.price)).scalar() or 0
        
        # Доходы
        total_revenue = session.query(func.sum(Order.price)).filter(Order.status.in_(["paid", "completed"])).scalar() or 0
        week_revenue = session.query(func.sum(Order.price)).filter(
            Order.status.in_(["paid", "completed"]),
            Order.created_at > week_ago
        ).scalar() or 0
        
        # Популярные типы работ
        popular_works = session.query(
            Order.work_type,
            func.count(Order.id).label('count')
        ).group_by(Order.work_type).order_by(func.count(Order.id).desc()).limit(3).all()
        
        # Популярные предметы
        popular_subjects = session.query(
            Order.science_name,
            func.count(Order.id).label('count')
        ).group_by(Order.science_name).order_by(func.count(Order.id).desc()).limit(3).all()
        
        # Среднее время выполнения заказа (для завершенных заказов)
        completed_orders_with_time = session.query(Order).filter(
            Order.status == "completed",
            Order.completed_at.isnot(None)
        ).all()
        
        if completed_orders_with_time:
            completion_times = [
                (order.completed_at - order.created_at).total_seconds() / 3600  # в часах
                for order in completed_orders_with_time
            ]
            avg_completion_time = sum(completion_times) / len(completion_times)
        else:
            avg_completion_time = 0
        
        # Активность по часам (топ-3 часа)
        from sqlalchemy import extract
        hourly_activity = session.query(
            extract('hour', UserAction.timestamp).label('hour'),
            func.count(UserAction.id).label('count')
        ).filter(UserAction.timestamp > day_ago).group_by(
            extract('hour', UserAction.timestamp)
        ).order_by(func.count(UserAction.id).desc()).limit(3).all()
        
        # Retention: пользователи, вернувшиеся после первого дня
        first_actions = session.query(
            UserAction.user_id,
            func.min(UserAction.timestamp).label('first_action')
        ).group_by(UserAction.user_id).subquery()
        
        returning_users = session.query(UserAction.user_id).join(
            first_actions, UserAction.user_id == first_actions.c.user_id
        ).filter(
            UserAction.timestamp > first_actions.c.first_action + timedelta(days=1)
        ).distinct().count()
        
        retention_rate = (returning_users / unique_users * 100) if unique_users > 0 else 0
        
        stats_text = f"""
📊 **Статистика бота:**

👥 Уникальных пользователей: {unique_users}
📝 Всего действий: {total_actions}
🚀 Команд /start: {start_commands}
📋 Команд /order: {order_commands}
📈 Среднее действий/пользователь: {avg_actions_per_user:.1f}
🔄 Retention rate: {retention_rate:.1f}%

💰 **Статистика заказов:**
📦 Всего заказов: {total_orders}
✅ Оплаченных: {paid_orders}
🎯 Завершенных: {completed_orders}
❌ Неудачных: {failed_orders}
📈 Конверсия: {conversion_rate:.1f}%
💵 Средняя цена: {avg_order_price:.0f} руб.
💸 Общий доход: {total_revenue:.0f} руб.

⏰ **Временная аналитика:**
🕐 Действий за 24 часа: {recent_actions}
📅 Действий за неделю: {week_actions}
🆕 Новых пользователей за неделю: {new_users_week}
🔥 Активных пользователей за неделю: {active_users_week}
📊 Активность сегодня: {today_actions}
🛒 Заказов за 24 часа: {recent_orders}
📋 Заказов за неделю: {week_orders}
💰 Доход за неделю: {week_revenue:.0f} руб.
⏱️ Среднее время выполнения: {avg_completion_time:.1f} ч.

🎯 **Популярные действия:**"""
        
        for action, count in top_actions:
            stats_text += f"\n• {action}: {count}"
        
        stats_text += f"""

📊 **Популярные типы работ:**"""
        for work_type, count in popular_works:
            stats_text += f"\n• {work_type}: {count}"
            
        stats_text += f"""

🎓 **Популярные предметы:**"""
        for subject, count in popular_subjects:
            stats_text += f"\n• {subject}: {count}"
            
        stats_text += f"""

🕐 **Активность по часам (топ-3):**"""
        for hour, count in hourly_activity:
            stats_text += f"\n• {int(hour)}:00 - {count} действий"
        
        stats_text += f"""

⚡ **Система:**
🔒 Администраторов: {len(ADMIN_IDS)}
🛡️ Rate limiting: {MAX_REQUESTS_PER_HOUR} req/hour
🎯 Семафор генерации: 10 слотов (свободно: {GENERATION_SEMAPHORE._value})
🕐 Rate limit записей: {len(user_request_times)}
        """
        
        await update.message.reply_text(stats_text, parse_mode='Markdown')
        
    except Exception as e:
        logging.error(f"Ошибка получения статистики: {e}")
        await update.message.reply_text("❌ Ошибка получения статистики")
    finally:
        session.close()

async def admin_users(update: Update, context: CallbackContext) -> None:
    """Список активных пользователей для администраторов"""
    if not is_admin(update.effective_user.id):
        await update.message.reply_text("❌ Доступ запрещен")
        return
    
    try:
        session = SessionLocal()
        
        # Получаем топ пользователей по активности
        from sqlalchemy import func, desc
        top_users = session.query(
            UserAction.user_id,
            func.count(UserAction.id).label('action_count'),
            func.max(UserAction.timestamp).label('last_activity')
        ).group_by(UserAction.user_id)\
         .order_by(desc('action_count'))\
         .limit(10).all()
        
        users_text = "👥 **Топ-10 активных пользователей:**\n\n"
        for user_id, action_count, last_activity in top_users:
            users_text += f"• User ID: {user_id}\n"
            users_text += f"  📝 Действий: {action_count}\n"
            users_text += f"  🕐 Последняя активность: {last_activity.strftime('%Y-%m-%d %H:%M')}\n\n"
        
        await update.message.reply_text(users_text, parse_mode='Markdown')
        
    except Exception as e:
        logging.error(f"Ошибка получения пользователей: {e}")
        await update.message.reply_text("❌ Ошибка получения данных пользователей")
    finally:
        session.close()

async def admin_broadcast(update: Update, context: CallbackContext) -> None:
    """Рассылка сообщений всем пользователям"""
    if not is_admin(update.effective_user.id):
        await update.message.reply_text("❌ Доступ запрещен")
        return
    
    if not context.args:
        await update.message.reply_text(
            "📢 Использование: /broadcast <сообщение>\n\n"
            "Пример: /broadcast Привет! У нас новые возможности!"
        )
        return
    
    message = ' '.join(context.args)
    
    try:
        session = SessionLocal()
        
        # Получаем уникальных пользователей
        unique_users = session.query(UserAction.user_id).distinct().all()
        user_ids = [user[0] for user in unique_users]
        
        success_count = 0
        error_count = 0
        
        await update.message.reply_text(f"📤 Начинаю рассылку для {len(user_ids)} пользователей...")
        
        for user_id in user_ids:
            try:
                await context.bot.send_message(chat_id=int(user_id), text=message)
                success_count += 1
                
                # Небольшая пауза для соблюдения лимитов Telegram API
                await asyncio.sleep(0.1)
                
            except Exception as e:
                error_count += 1
                logging.warning(f"Не удалось отправить сообщение пользователю {user_id}: {e}")
        
        result_text = f"""
✅ **Рассылка завершена!**

📤 Отправлено: {success_count}
❌ Ошибок: {error_count}
📊 Всего попыток: {len(user_ids)}
        """
        
        await update.message.reply_text(result_text, parse_mode='Markdown')
        
    except Exception as e:
        logging.error(f"Ошибка рассылки: {e}")
        await update.message.reply_text("❌ Ошибка при выполнении рассылки")
    finally:
        session.close()

async def admin_system(update: Update, context: CallbackContext) -> None:
    """Информация о системе"""
    if not is_admin(update.effective_user.id):
        await update.message.reply_text("❌ Доступ запрещен")
        return
    
    try:
        import psutil
        import sys
        
        # Системная информация
        cpu_percent = psutil.cpu_percent(interval=1)
        memory = psutil.virtual_memory()
        disk = psutil.disk_usage('/')
        
        # Информация о Python процессе
        process = psutil.Process()
        bot_memory = process.memory_info().rss / 1024 / 1024  # MB
        
        system_text = f"""
🖥️ **Системная информация:**

💻 CPU: {cpu_percent}%
🧠 RAM: {memory.percent}% ({memory.used // 1024 // 1024} MB / {memory.total // 1024 // 1024} MB)
💾 Диск: {disk.percent}% ({disk.used // 1024 // 1024 // 1024} GB / {disk.total // 1024 // 1024 // 1024} GB)

🤖 **Бот:**
📊 Память бота: {bot_memory:.1f} MB
🐍 Python: {sys.version.split()[0]}
⚡ Активных генераций: {10 - GENERATION_SEMAPHORE._value}
🕐 Rate limit записей: {len(user_request_times)}
        """
        
        await update.message.reply_text(system_text, parse_mode='Markdown')
        
    except ImportError:
        await update.message.reply_text(
            "📋 **Базовая информация:**\n\n"
            "⚡ Семафор генерации: доступно слотов " + str(GENERATION_SEMAPHORE._value) + "/10\n"
            "🕐 Rate limit записей: " + str(len(user_request_times)) + "\n\n"
            "ℹ️ Для полной информации установите: pip install psutil"
        )
    except Exception as e:
        logging.error(f"Ошибка получения системной информации: {e}")
        await update.message.reply_text("❌ Ошибка получения системной информации")

async def admin_orders(update: Update, context: CallbackContext) -> None:
    """Просмотр заказов для администраторов"""
    if not is_admin(update.effective_user.id):
        await update.message.reply_text("❌ Доступ запрещен")
        return
    
    try:
        session = SessionLocal()
        
        # Статистика по заказам
        from sqlalchemy import func, desc
        
        total_orders = session.query(Order).count()
        completed_orders = session.query(Order).filter(Order.status == "completed").count()
        failed_orders = session.query(Order).filter(Order.status == "failed").count()
        refunded_orders = session.query(Order).filter(Order.status == "refunded").count()
        
        # Последние заказы
        recent_orders = session.query(Order)\
            .order_by(desc(Order.created_at))\
            .limit(5).all()
        
        # Доходы
        total_revenue = session.query(func.sum(Order.price))\
            .filter(Order.status == "completed").scalar() or 0
        
        orders_text = f"""
📊 **Статистика заказов:**

📝 Всего заказов: {total_orders}
✅ Выполнено: {completed_orders}
❌ Ошибки: {failed_orders}
💸 Возвраты: {refunded_orders}
💰 Общий доход: {total_revenue}₽

📋 **Последние заказы:**

        """
        
        for order in recent_orders:
            status_emoji = {
                "created": "🆕",
                "paid": "💳",
                "completed": "✅",
                "failed": "❌",
                "refunded": "💸",
                "cancelled": "🚫"
            }.get(order.status, "❓")
            
            orders_text += f"""
{status_emoji} ID: {order.id} | {order.work_type}
💰 {order.price}₽ | 👤 {order.user_id}
📚 {order.science_name} - {order.work_theme[:30]}...
🕐 {order.created_at.strftime('%d.%m %H:%M')}

"""
        
        await update.message.reply_text(orders_text, parse_mode='Markdown')
        
    except Exception as e:
        logging.error(f"Ошибка получения заказов: {e}")
        await update.message.reply_text("❌ Ошибка получения данных заказов")
    finally:
        session.close()

async def admin_finance(update: Update, context: CallbackContext) -> None:
    """Детальная финансовая аналитика для администраторов"""
    if not is_admin(update.effective_user.id):
        await update.message.reply_text("❌ Доступ запрещен")
        return
    
    try:
        session = SessionLocal()
        from sqlalchemy import func, desc
        
        # Временные интервалы
        now = datetime.now(timezone.utc)
        today = now.replace(hour=0, minute=0, second=0, microsecond=0)
        yesterday = today - timedelta(days=1)
        week_ago = now - timedelta(days=7)
        month_ago = now - timedelta(days=30)
        
        # Доходы по периодам
        today_revenue = session.query(func.sum(Order.price)).filter(
            Order.status.in_(["paid", "completed"]),
            Order.created_at >= today
        ).scalar() or 0
        
        yesterday_revenue = session.query(func.sum(Order.price)).filter(
            Order.status.in_(["paid", "completed"]),
            Order.created_at >= yesterday,
            Order.created_at < today
        ).scalar() or 0
        
        week_revenue = session.query(func.sum(Order.price)).filter(
            Order.status.in_(["paid", "completed"]),
            Order.created_at >= week_ago
        ).scalar() or 0
        
        month_revenue = session.query(func.sum(Order.price)).filter(
            Order.status.in_(["paid", "completed"]),
            Order.created_at >= month_ago
        ).scalar() or 0
        
        # Заказы по периодам
        today_orders = session.query(Order).filter(Order.created_at >= today).count()
        yesterday_orders = session.query(Order).filter(
            Order.created_at >= yesterday,
            Order.created_at < today
        ).count()
        
        # Средний чек по периодам
        today_avg = today_revenue / today_orders if today_orders > 0 else 0
        yesterday_avg = yesterday_revenue / yesterday_orders if yesterday_orders > 0 else 0
        
        # Топ-5 самых дорогих заказов
        expensive_orders = session.query(Order).filter(
            Order.status.in_(["paid", "completed"])
        ).order_by(desc(Order.price)).limit(5).all()
        
        # Статистика по типам работ (доходы)
        work_revenue = session.query(
            Order.work_type,
            func.sum(Order.price).label('revenue'),
            func.count(Order.id).label('count')
        ).filter(
            Order.status.in_(["paid", "completed"])
        ).group_by(Order.work_type).order_by(desc('revenue')).limit(5).all()
        
        # Статистика по дням недели
        weekday_stats = session.query(
            func.extract('dow', Order.created_at).label('weekday'),
            func.sum(Order.price).label('revenue'),
            func.count(Order.id).label('count')
        ).filter(
            Order.status.in_(["paid", "completed"]),
            Order.created_at >= week_ago
        ).group_by(func.extract('dow', Order.created_at)).all()
        
        # Конверсия воронки
        total_users = session.query(UserAction.user_id).distinct().count()
        order_users = session.query(Order.user_id).distinct().count()
        paid_users = session.query(Order.user_id).filter(
            Order.status.in_(["paid", "completed"])
        ).distinct().count()
        
        conversion_to_order = (order_users / total_users * 100) if total_users > 0 else 0
        conversion_to_payment = (paid_users / order_users * 100) if order_users > 0 else 0
        
        # Возвраты и отмены
        refunds = session.query(Order).filter(Order.status == "refunded").count()
        failed_orders = session.query(Order).filter(Order.status == "failed").count()
        
        finance_text = f"""
💰 **Финансовая аналитика:**

📊 **Доходы:**
💸 Сегодня: {today_revenue:.0f} руб. ({today_orders} заказов)
📅 Вчера: {yesterday_revenue:.0f} руб. ({yesterday_orders} заказов)
📈 За неделю: {week_revenue:.0f} руб.
📊 За месяц: {month_revenue:.0f} руб.

💵 **Средний чек:**
🔥 Сегодня: {today_avg:.0f} руб.
📊 Вчера: {yesterday_avg:.0f} руб.

🎯 **Конверсия:**
👥 Всего пользователей: {total_users}
📝 Создали заказ: {order_users} ({conversion_to_order:.1f}%)
💳 Оплатили: {paid_users} ({conversion_to_payment:.1f}%)

🔥 **Топ-5 дорогих заказов:**"""
        
        for order in expensive_orders:
            finance_text += f"\n💰 {order.price}₽ - {order.work_type} ({order.science_name})"
        
        finance_text += f"""

📊 **Доходы по типам работ:**"""
        for work_type, revenue, count in work_revenue:
            finance_text += f"\n• {work_type}: {revenue:.0f}₽ ({count} заказов)"
        
        finance_text += f"""

📅 **Статистика по дням недели:**"""
        weekdays = ['Пн', 'Вт', 'Ср', 'Чт', 'Пт', 'Сб', 'Вс']
        for weekday, revenue, count in weekday_stats:
            day_name = weekdays[int(weekday) - 1] if weekday and weekday > 0 else "Неизв"
            finance_text += f"\n• {day_name}: {revenue:.0f}₽ ({count} заказов)"
        
        finance_text += f"""

⚠️ **Проблемы:**
🔴 Возвраты: {refunds}
❌ Неудачные заказы: {failed_orders}
        """
        
        await update.message.reply_text(finance_text, parse_mode='Markdown')
        
    except Exception as e:
        logging.error(f"Ошибка финансовой аналитики: {e}")
        await update.message.reply_text("❌ Ошибка получения финансовой аналитики")
    finally:
        session.close()

async def admin_export(update: Update, context: CallbackContext) -> None:
    """Экспорт данных в CSV для администраторов"""
    if not is_admin(update.effective_user.id):
        await update.message.reply_text("❌ Доступ запрещен")
        return
    
    try:
        session = SessionLocal()
        import csv
        from io import StringIO
        
        # Экспорт заказов
        orders = session.query(Order).all()
        
        # Создаем CSV в памяти
        csv_buffer = StringIO()
        writer = csv.writer(csv_buffer)
        
        # Заголовки
        writer.writerow([
            'ID', 'User ID', 'Work Type', 'Science Name', 'Theme', 
            'Pages', 'Price', 'Status', 'Created At', 'Completed At'
        ])
        
        # Данные
        for order in orders:
            writer.writerow([
                order.id,
                order.user_id,
                order.work_type,
                order.science_name,
                order.work_theme,
                order.page_number,
                order.price,
                order.status,
                order.created_at.strftime('%Y-%m-%d %H:%M:%S'),
                order.completed_at.strftime('%Y-%m-%d %H:%M:%S') if order.completed_at else ''
            ])
        
        # Отправляем файл
        csv_content = csv_buffer.getvalue().encode('utf-8')
        await update.message.reply_document(
            document=io.BytesIO(csv_content),
            filename=f"orders_export_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv",
            caption="📊 Экспорт заказов в CSV"
        )
        
        # Экспорт активности пользователей
        actions = session.query(UserAction).all()
        
        actions_csv_buffer = StringIO()
        actions_writer = csv.writer(actions_csv_buffer)
        
        # Заголовки
        actions_writer.writerow(['ID', 'User ID', 'Action', 'Timestamp'])
        
        # Данные
        for action in actions:
            actions_writer.writerow([
                action.id,
                action.user_id,
                action.action,
                action.timestamp.strftime('%Y-%m-%d %H:%M:%S')
            ])
        
        # Отправляем файл активности
        actions_csv_content = actions_csv_buffer.getvalue().encode('utf-8')
        await update.message.reply_document(
            document=io.BytesIO(actions_csv_content),
            filename=f"user_actions_export_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv",
            caption="📊 Экспорт активности пользователей в CSV"
        )
        
        await update.message.reply_text("✅ Экспорт данных завершен!")
        
    except Exception as e:
        logging.error(f"Ошибка экспорта данных: {e}")
        await update.message.reply_text("❌ Ошибка экспорта данных")
    finally:
        session.close()

async def admin_monitor(update: Update, context: CallbackContext) -> None:
    """Мониторинг системы в реальном времени"""
    if not is_admin(update.effective_user.id):
        await update.message.reply_text("❌ Доступ запрещен")
        return
    
    try:
        session = SessionLocal()
        
        # Активность за последние 5 минут
        last_5_min = datetime.now(timezone.utc) - timedelta(minutes=5)
        recent_activity = session.query(UserAction).filter(
            UserAction.timestamp > last_5_min
        ).count()
        
        # Активные генерации
        active_generations = 10 - GENERATION_SEMAPHORE._value
        
        # Заказы в очереди (created статус)
        pending_orders = session.query(Order).filter(Order.status == "created").count()
        
        # Неудачные заказы за последний час
        last_hour = datetime.now(timezone.utc) - timedelta(hours=1)
        failed_recent = session.query(Order).filter(
            Order.status == "failed",
            Order.created_at > last_hour
        ).count()
        
        # Системные метрики
        try:
            import psutil
            cpu_percent = psutil.cpu_percent(interval=1)
            memory_percent = psutil.virtual_memory().percent
            system_status = f"💻 CPU: {cpu_percent}% | 🧠 RAM: {memory_percent}%"
        except ImportError:
            system_status = "📊 Системные метрики недоступны"
        
        # Статус семафора
        semaphore_status = "🟢 Норма" if active_generations < 8 else "🟡 Высокая нагрузка" if active_generations < 10 else "🔴 Перегрузка"
        
        # Алерты
        alerts = []
        if failed_recent > 5:
            alerts.append("🚨 Много неудачных заказов за час")
        if active_generations >= 9:
            alerts.append("⚠️ Высокая нагрузка генерации")
        if pending_orders > 10:
            alerts.append("📋 Много заказов в очереди")
        
        monitor_text = f"""
🔍 **Мониторинг системы:**

⚡ **Активность:**
🕐 Действий за 5 мин: {recent_activity}
🎯 Активных генераций: {active_generations}/10
📋 Заказов в очереди: {pending_orders}
❌ Неудачных заказов/час: {failed_recent}

🖥️ **Система:**
{system_status}
🎯 Статус генерации: {semaphore_status}

🚨 **Алерты:**"""
        
        if alerts:
            for alert in alerts:
                monitor_text += f"\n{alert}"
        else:
            monitor_text += "\n✅ Все системы работают нормально"
        
        monitor_text += f"""

📊 **Быстрые действия:**
/admin_stats - Полная статистика
/admin_finance - Финансовая аналитика
/admin_system - Системная информация
        """
        
        await update.message.reply_text(monitor_text, parse_mode='Markdown')
        
    except Exception as e:
        logging.error(f"Ошибка мониторинга: {e}")
        await update.message.reply_text("❌ Ошибка получения данных мониторинга")
    finally:
        session.close()

# Настройка базы данных
DATABASE_URL = "sqlite:///user_activity.db"
Base = declarative_base()  # Updated to use sqlalchemy.orm.declarative_base
engine = create_engine(DATABASE_URL)
SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)

# Модель для хранения действий пользователей
class UserAction(Base):
    __tablename__ = "user_actions"

    id = Column(Integer, primary_key=True, index=True)
    user_id = Column(String, index=True)
    action = Column(String, index=True)
    timestamp = Column(DateTime, default=lambda: datetime.now(timezone.utc))

# Новая модель для хранения заказов
class Order(Base):
    __tablename__ = "orders"
    
    id = Column(Integer, primary_key=True, index=True)
    user_id = Column(String, index=True)
    work_type = Column(String)
    science_name = Column(String)
    work_theme = Column(String)
    page_number = Column(Integer)
    price = Column(Integer)
    status = Column(String, default="created")  # created, paid, completed, failed, refunded
    payment_id = Column(String, nullable=True)
    created_at = Column(DateTime, default=lambda: datetime.now(timezone.utc))
    completed_at = Column(DateTime, nullable=True)

# Создание таблиц
Base.metadata.create_all(bind=engine)

# Функция для записи действий пользователя
async def log_user_action(user_id: str, action: str):
    session = SessionLocal()
    try:
        user_action = UserAction(
            user_id=str(user_id), 
            action=action,
            timestamp=datetime.now(timezone.utc)  # Заменить deprecated utcnow()
        )
        session.add(user_action)
        session.commit()
        logging.info(f"User action logged: {action} for user {user_id}")
    except Exception as e:
        session.rollback()
        logging.error(f"Ошибка записи действия пользователя: {e}")
        raise
    finally:
        session.close()

# Функция для создания заказа
async def create_order(user_id: str, order_data: dict) -> int:
    """Создает заказ в базе данных и возвращает ID заказа"""
    session = SessionLocal()
    try:
        order = Order(
            user_id=str(user_id),
            work_type=order_data.get("work_type", ""),
            science_name=order_data.get("science_name", ""),
            work_theme=order_data.get("work_theme", ""),
            page_number=order_data.get("page_number", 0),
            price=order_data.get("price", 0),
            status="created"
        )
        session.add(order)
        session.commit()
        order_id = order.id
        logging.info(f"Order created: {order_id} for user {user_id}")
        return order_id
    except Exception as e:
        session.rollback()
        logging.error(f"Ошибка создания заказа: {e}")
        raise
    finally:
        session.close()

# Функция для обновления статуса заказа
async def update_order_status(order_id: int, status: str, payment_id: str = None):
    """Обновляет статус заказа"""
    session = SessionLocal()
    try:
        order = session.query(Order).filter(Order.id == order_id).first()
        if order:
            order.status = status
            if payment_id:
                order.payment_id = payment_id
            if status == "completed":
                order.completed_at = datetime.now(timezone.utc)
            session.commit()
            logging.info(f"Order {order_id} status updated to {status}")
    except Exception as e:
        session.rollback()
        logging.error(f"Ошибка обновления статуса заказа: {e}")
        raise
    finally:
        session.close()

# Команда /start
async def start(update: Update, context: CallbackContext) -> None:
    user_id = update.effective_user.id
    username = update.effective_user.username or ""
    await log_user_action(user_id, "start_command")
    # Show user agreement before displaying menu
    keyboard = [["Продолжить"]]
    reply_markup = ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
    
    # Уведомление о тестовом режиме
    mode_text = ""
    is_test_user = username in TEST_MODE_USERNAMES
    if TESTING_MODE or is_test_user:
        mode_text = "\n🧪 *ВНИМАНИЕ: БОТ В ТЕСТОВОМ РЕЖИМЕ* 🧪\n" \
                   "Все заказы выполняются БЕСПЛАТНО для тестирования!\n\n"
    
    await update.message.reply_text(
        "👾 Добро пожаловать в *NinjaEssayAI*! 🥷\n\n"
        f"{mode_text}"
        "Перед началом работы, пожалуйста, ознакомьтесь с пользовательским соглашением:\n"
        "https://docs.google.com/document/d/100ljVD3fFveH3Vuz7Y6F9QfFj5EsdOBymdRDRAMe2MI/edit?tab=t.0\n\n"
        "Нажимая на кнопку \"Продолжить\", вы подтверждаете, что ознакомились с пользовательским соглашением.",
        reply_markup=reply_markup,
        parse_mode='Markdown'
    )

# Команда /help
async def help_command(update: Update, context: CallbackContext) -> None:
    await update.message.reply_text(
        "🆘 *Помощь* 🆘\n\n"
        "Вот список доступных команд:\n"
        "- 📋 /order - Оформить заказ\n"
        "- ❌ /cancel - Отменить текущий процесс\n"
        "- ℹ️ /help - Показать это сообщение помощи"
    )

# Добавление команды /menu для отображения меню в любой момент диалога
async def menu(update: Update, context: CallbackContext) -> None:
    keyboard = [["/order"], ["/help", "/cancel"]]
    reply_markup = ReplyKeyboardMarkup(keyboard, resize_keyboard=True)
    await update.message.reply_text(
        "📋 *Меню команд:* 📋\n\n"
        "- /order - Оформить заказ\n"
        "- /help - Помощь\n"
        "- /cancel - Отмена текущего процесса",
        reply_markup=reply_markup
    )

# Добавляем обработчик принятия пользовательского соглашения
async def continue_handler(update: Update, context: CallbackContext) -> None:
    user_id = update.effective_user.id
    await log_user_action(user_id, "accepted_terms")
    keyboard = [["/order"], ["/help", "/cancel"]]
    reply_markup = ReplyKeyboardMarkup(keyboard, resize_keyboard=True)
    await update.message.reply_text(
        "📋 *Меню команд:* 📋\n\n"
        "- /order - Оформить заказ\n"
        "- /help - Помощь\n"
        "- /cancel - Отмена текущего процесса",
        reply_markup=reply_markup
    )

# Состояния разговора
WORK_TYPE, SCIENCE_NAME, PAGE_NUMBER, WORK_THEME, CUSTOM_PLAN, PREFERENCES, PAYMENT = range(7)

# Функция для создания клавиатуры с кнопкой "Назад"
def create_keyboard_with_back(options, show_back=True):
    """Создает клавиатуру с опциональной кнопкой 'Назад'"""
    keyboard = []
    
    # Добавляем основные опции только если они есть
    if options:  # Проверяем, что список не пустой
        if isinstance(options[0], list):
            keyboard.extend(options)
        else:
            # Если это простой список строк, группируем по 2
            for i in range(0, len(options), 2):
                row = options[i:i+2]
                keyboard.append(row)
    
    # Добавляем кнопку "Назад" если нужно
    if show_back:
        keyboard.append(["◀️ Назад"])
    
    return keyboard

# Функция для возврата на предыдущий шаг
async def go_back_handler(update: Update, context: CallbackContext) -> int:
    """Обработчик возврата на предыдущий шаг"""
    current_step = context.user_data.get("current_step", WORK_TYPE)
    
    logging.info(f"Пользователь нажал 'Назад' на шаге: {current_step}")
    
    if current_step == WORK_TYPE:
        await update.message.reply_text("Вы уже на первом шаге. Используйте /cancel для отмены заказа.")
        return await order(update, context)
    elif current_step == SCIENCE_NAME:
        logging.info("Возврат к выбору типа работы")
        return await science_name_back(update, context)
    elif current_step == PAGE_NUMBER:
        logging.info("Возврат к вводу дисциплины")
        return await page_number_back(update, context)
    elif current_step == WORK_THEME:
        logging.info("Возврат к вводу количества страниц")
        return await work_theme_back(update, context)
    elif current_step == CUSTOM_PLAN:
        logging.info("Возврат к вводу темы работы")
        return await custom_plan_back(update, context)
    elif current_step == PREFERENCES:
        logging.info("Возврат к выбору типа плана")
        return await preferences_back(update, context)
    elif current_step == PAYMENT:
        logging.info("Возврат к предпочтениям (с этапа оплаты)")
        return await preferences_back(update, context)
    
    return WORK_TYPE

# Функции для возврата на конкретные шаги

async def science_name_back(update: Update, context: CallbackContext) -> int:
    """Возврат к выбору типа работы"""
    keyboard = create_keyboard_with_back([
        ["📝 Эссе - 300₽", "📜 Доклад - 300₽"],
        ["📖 Реферат - 400₽", "💼 Проект - 400₽"],
        ["📚 Курсовая работа - 500₽"],
        ["🎓 Дипломная работа - 800₽"]
    ], show_back=False)
    reply_markup = ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
    await update.message.reply_text("🥷 *Выберите тип работы:* 🥷", reply_markup=reply_markup)
    context.user_data["current_step"] = WORK_TYPE
    return WORK_TYPE

async def page_number_back(update: Update, context: CallbackContext) -> int:
    """Возврат к вводу дисциплины"""
    keyboard = create_keyboard_with_back([], show_back=True)
    reply_markup = ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
    await update.message.reply_text(
        "📝 Шаг 1/6: Укажите название дисциплины (например, Математика):",
        reply_markup=reply_markup
    )
    context.user_data["current_step"] = SCIENCE_NAME
    return SCIENCE_NAME

async def work_theme_back(update: Update, context: CallbackContext) -> int:
    """Возврат к вводу количества страниц"""
    work_type = context.user_data.get("work_type", "")
    clean_type = re.sub(r"[^А-Яа-яЁё ]", "", work_type).strip()
    max_pages = PAGE_LIMITS.get(clean_type, 10)
    
    keyboard = create_keyboard_with_back([], show_back=True)
    reply_markup = ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
    await update.message.reply_text(
        f"📝 Шаг 2/6: Введите количество страниц (максимум {max_pages}):",
        reply_markup=reply_markup
    )
    context.user_data["current_step"] = PAGE_NUMBER
    return PAGE_NUMBER

async def custom_plan_back(update: Update, context: CallbackContext) -> int:
    """Возврат к вводу темы работы"""
    keyboard = create_keyboard_with_back([], show_back=True)
    reply_markup = ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
    await update.message.reply_text(
        "📝 Шаг 3/6: Укажите тему работы:",
        reply_markup=reply_markup
    )
    context.user_data["current_step"] = WORK_THEME
    # Сбрасываем данные о плане при возврате
    context.user_data.pop("plan_entered", None)
    context.user_data.pop("custom_plan", None)
    context.user_data.pop("use_custom_plan", None)
    return WORK_THEME

async def preferences_back(update: Update, context: CallbackContext) -> int:
    """Возврат к выбору типа плана"""
    # Сбрасываем данные о предпочтениях при возврате
    context.user_data.pop("preferences", None)
    
    # Проверяем, был ли введён пользовательский план
    use_custom_plan = context.user_data.get("use_custom_plan", False)
    plan_entered = context.user_data.get("plan_entered", False)
    
    if use_custom_plan and plan_entered:
        # Если план уже был введён, возвращаемся к выбору типа плана
        # но сбрасываем информацию о введённом плане
        context.user_data.pop("plan_entered", None)
        context.user_data.pop("custom_plan", None)
        context.user_data.pop("use_custom_plan", None)
        
        keyboard = create_keyboard_with_back([
            ["🤖 Автоматический план"],
            ["✍️ Создать свой план"]
        ], show_back=True)
        reply_markup = ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
        await update.message.reply_text(
            "📝 Шаг 4/6: Выберите способ создания плана работы:",
            reply_markup=reply_markup
        )
        context.user_data["current_step"] = CUSTOM_PLAN
        return CUSTOM_PLAN
    else:
        # Если план не был введён, также возвращаемся к выбору типа плана
        keyboard = create_keyboard_with_back([
            ["🤖 Автоматический план"],
            ["✍️ Создать свой план"]
        ], show_back=True)
        reply_markup = ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
        await update.message.reply_text(
            "📝 Шаг 4/6: Выберите способ создания плана работы:",
            reply_markup=reply_markup
        )
        context.user_data["current_step"] = CUSTOM_PLAN
        # Сбрасываем данные о плане при возврате
        context.user_data.pop("plan_entered", None)
        context.user_data.pop("custom_plan", None)
        context.user_data.pop("use_custom_plan", None)
        return CUSTOM_PLAN

# Ограничение страниц по типам работы
PAGE_LIMITS = {"Эссе": 10, "Доклад": 10, "Реферат": 20, "Проект": 20, "Курсовая работа": 30, "Дипломная работа": 70}
# Семафор для ограничения одновременных запросов к DeepSeek API
# Настройки семафора:
# 5 - для малой нагрузки (1-10 пользователей)
# 10 - для средней нагрузки (10-50 пользователей) [текущая]
# 20 - для высокой нагрузки (50+ пользователей)
# 30 - для максимальной нагрузки (100+ пользователей)
GENERATION_SEMAPHORE = asyncio.Semaphore(10)

# Начало заказа
async def order(update: Update, context: CallbackContext) -> int:
    user_id = update.effective_user.id
    
    await log_user_action(user_id, "order_command")
    # Updated keyboard to include prices for each work type
    keyboard = create_keyboard_with_back([
        ["📝 Эссе - 300₽", "📜 Доклад - 300₽"],
        ["📖 Реферат - 400₽", "💼 Проект - 400₽"],
        ["📚 Курсовая работа - 500₽"],
        ["🎓 Дипломная работа - 800₽"]
    ], show_back=False)
    reply_markup = ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
    await update.message.reply_text("🥷 *Выберите тип работы:* 🥷", reply_markup=reply_markup)
    context.user_data["current_step"] = WORK_TYPE
    return WORK_TYPE

# Обработчик inline кнопки "Назад"
async def back_button_handler(update: Update, context: CallbackContext) -> int:
    """Обработчик inline кнопки 'Назад' на этапе оплаты"""
    query = update.callback_query
    await query.answer()
    
    # Получаем chat_id для отправки сообщений
    chat_id = query.message.chat_id
    
    # Возвращаемся к этапу предпочтений
    use_custom_plan = context.user_data.get("use_custom_plan", False)
    plan_entered = context.user_data.get("plan_entered", False)
    
    # Сбрасываем данные о предпочтениях
    context.user_data.pop("preferences", None)
    context.user_data["current_step"] = PREFERENCES
    
    # Удаляем сообщение с кнопкой оплаты
    try:
        await query.message.delete()
    except:
        pass
    
    if use_custom_plan and plan_entered:
        # Если был пользовательский план, возвращаемся к вводу предпочтений
        keyboard = create_keyboard_with_back([["⏭️ Пропустить"]], show_back=True)
        reply_markup = ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
        await context.bot.send_message(
            chat_id=chat_id,
            text="📝 Шаг 6/6: Опишите ваши предпочтения по работе (например, стиль, источники, сроки):",
            reply_markup=reply_markup
        )
    elif use_custom_plan and not plan_entered:
        # Возвращаемся к вводу плана
        keyboard = create_keyboard_with_back([], show_back=True)
        reply_markup = ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
        await context.bot.send_message(
            chat_id=chat_id,
            text="📝 Шаг 5/6: Введите план вашей работы.\n\n"
                "Каждый пункт плана с новой строки, например:\n"
                "1. Введение\n"
                "2. Основная часть\n"
                "3. Заключение\n\n"
                "Или просто:\n"
                "Введение\n"
                "Основная часть\n"
                "Заключение",
            reply_markup=reply_markup
        )
    else:
        # Был автоматический план, возвращаемся к предпочтениям
        keyboard = create_keyboard_with_back([["⏭️ Пропустить"]], show_back=True)
        reply_markup = ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
        await context.bot.send_message(
            chat_id=chat_id,
            text="📝 Шаг 5/6: Опишите ваши предпочтения по работе (например, стиль, источники, сроки):",
            reply_markup=reply_markup
        )
    
    return PREFERENCES

# Обработчики этапов заказа
async def work_type_handler(update: Update, context: CallbackContext) -> int:
    # Проверяем, не нажал ли пользователь "Назад"
    if update.message.text == "◀️ Назад":
        return await go_back_handler(update, context)
    
    # Parse selected work type and price if present
    selection = update.message.text
    if " - " in selection:
        type_text, price_text = selection.split(" - ", 1)
        context.user_data["work_type"] = type_text
        try:
            context.user_data["price"] = int(price_text.replace("₽","").replace(" ", ""))
        except ValueError:
            context.user_data["price"] = None
    else:
        context.user_data["work_type"] = selection
        context.user_data["price"] = None
    
    # Шаг 1 из 6: дисциплина
    keyboard = create_keyboard_with_back([], show_back=True)
    reply_markup = ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
    await update.message.reply_text(
        "📝 Шаг 1/6: Укажите название дисциплины (например, Математика):",
        reply_markup=reply_markup
    )
    context.user_data["current_step"] = SCIENCE_NAME
    return SCIENCE_NAME

async def science_name_handler(update: Update, context: CallbackContext) -> int:
    # Проверяем, не нажал ли пользователь "Назад"
    if update.message.text == "◀️ Назад":
        return await go_back_handler(update, context)
    
    try:
        science_name = validate_user_input(update.message.text, 100)
        context.user_data["science_name"] = science_name
    except ValueError as e:
        await update.message.reply_text(f"⚠️ Ошибка: {e}")
        return SCIENCE_NAME
    
    # Шаг 2 из 6: количество страниц (с учётом лимитов)
    work_type = context.user_data.get("work_type", "")
    clean_type = re.sub(r"[^А-Яа-яЁё ]", "", work_type).strip()
    max_pages = PAGE_LIMITS.get(clean_type, 10)
    
    keyboard = create_keyboard_with_back([], show_back=True)
    reply_markup = ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
    await update.message.reply_text(
        f"📝 Шаг 2/6: Введите количество страниц (максимум {max_pages}):",
        reply_markup=reply_markup
    )
    context.user_data["current_step"] = PAGE_NUMBER
    return PAGE_NUMBER

async def page_number_handler(update: Update, context: CallbackContext) -> int:
    # Проверяем, не нажал ли пользователь "Назад"
    if update.message.text == "◀️ Назад":
        return await go_back_handler(update, context)
    
    # Валидация числа страниц
    try:
        page = int(update.message.text)
    except ValueError:
        await update.message.reply_text("⚠️ Пожалуйста, введите число:")
        return PAGE_NUMBER
    work_type = context.user_data.get("work_type", "")
    clean_type = re.sub(r"[^А-Яа-яЁё ]", "", work_type).strip()
    max_pages = PAGE_LIMITS.get(clean_type, 10)
    if page < 1 or page > max_pages:
        await update.message.reply_text(
            f"❗️Для {clean_type} максимальное количество страниц — {max_pages}. Пожалуйста, введите число от 1 до {max_pages}:"
        )
        return PAGE_NUMBER
    context.user_data["page_number"] = page
    
    # Шаг 3 из 6: тема работы
    keyboard = create_keyboard_with_back([], show_back=True)
    reply_markup = ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
    await update.message.reply_text(
        "📝 Шаг 3/6: Укажите тему работы:",
        reply_markup=reply_markup
    )
    context.user_data["current_step"] = WORK_THEME
    return WORK_THEME

async def work_theme_handler(update: Update, context: CallbackContext) -> int:
    # Проверяем, не нажал ли пользователь "Назад"
    if update.message.text == "◀️ Назад":
        return await go_back_handler(update, context)
    
    try:
        work_theme = validate_user_input(update.message.text, 400)
        context.user_data["work_theme"] = work_theme
    except ValueError as e:
        await update.message.reply_text(f"⚠️ Ошибка: {e}")
        return WORK_THEME
        
    # Предлагаем выбор: автоматический план или пользовательский
    keyboard = create_keyboard_with_back([
        ["🤖 Автоматический план"],
        ["✍️ Создать свой план"]
    ], show_back=True)
    reply_markup = ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
    await update.message.reply_text(
        "📝 Шаг 4/6: Выберите способ создания плана работы:",
        reply_markup=reply_markup
    )
    context.user_data["current_step"] = CUSTOM_PLAN
    return CUSTOM_PLAN

async def custom_plan_handler(update: Update, context: CallbackContext) -> int:
    # Проверяем, не нажал ли пользователь "Назад"
    if update.message.text == "◀️ Назад":
        return await go_back_handler(update, context)
    
    choice = update.message.text
    
    if choice == "🤖 Автоматический план":
        # Автоматический план - переходим к предпочтениям
        context.user_data["use_custom_plan"] = False
        keyboard = create_keyboard_with_back([["⏭️ Пропустить"]], show_back=True)
        reply_markup = ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
        await update.message.reply_text(
            "📝 Шаг 5/6: Опишите ваши предпочтения по работе (например, стиль, источники, сроки):",
            reply_markup=reply_markup
        )
        context.user_data["current_step"] = PREFERENCES
        return PREFERENCES
    elif choice == "✍️ Создать свой план":
        # Пользовательский план
        context.user_data["use_custom_plan"] = True
        keyboard = create_keyboard_with_back([], show_back=True)
        reply_markup = ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
        await update.message.reply_text(
            "📝 Шаг 5/6: Введите план вашей работы.\n\n"
            "Каждый пункт плана с новой строки, например:\n"
            "1. Введение\n"
            "2. Основная часть\n"
            "3. Заключение\n\n"
            "Или просто:\n"
            "Введение\n"
            "Основная часть\n"
            "Заключение",
            reply_markup=reply_markup
        )
        context.user_data["current_step"] = PREFERENCES
        return PREFERENCES
    else:
        await update.message.reply_text(
            "⚠️ Пожалуйста, выберите один из предложенных вариантов:",
        )
        return CUSTOM_PLAN

async def preferences_handler(update: Update, context: CallbackContext) -> int:
    # Проверяем, не нажал ли пользователь "Назад"
    if update.message.text == "◀️ Назад":
        return await go_back_handler(update, context)
    
    use_custom_plan = context.user_data.get("use_custom_plan", False)
    plan_entered = context.user_data.get("plan_entered", False)
    
    if use_custom_plan and not plan_entered:
        # Обрабатываем пользовательский план
        try:
            custom_plan_text = validate_user_input(update.message.text, 1000)
            # Разбиваем план на отдельные пункты
            plan_lines = [line.strip() for line in custom_plan_text.split('\n') if line.strip()]
            
            # Очищаем от нумерации, если есть
            cleaned_plan = []
            for line in plan_lines:
                # Убираем нумерацию в начале строки
                cleaned_line = re.sub(r'^\d+\.?\s*', '', line).strip()
                if cleaned_line:
                    cleaned_plan.append(cleaned_line)
            
            if len(cleaned_plan) < 2:
                await update.message.reply_text(
                    "⚠️ План должен содержать минимум 2 пункта. Попробуйте еще раз:"
                )
                return PREFERENCES
            
            # Проверяем соответствие количества пунктов заявленному количеству страниц
            page_number = context.user_data.get("page_number", 0)
            expected_chapters = max(1, page_number // 2)  # Минимум 1 пункт
            plan_chapters = len(cleaned_plan)
            
            warning_text = ""
            if plan_chapters < expected_chapters:
                warning_text = f"\n\n⚠️ Рекомендация: Для {page_number} страниц обычно требуется {expected_chapters} пунктов плана, " \
                              f"а у вас {plan_chapters}. Это может привести к меньшему объему итоговой работы."
            elif plan_chapters > expected_chapters * 1.5:  # Если пунктов слишком много
                warning_text = f"\n\n⚠️ Внимание: У вас {plan_chapters} пунктов плана для {page_number} страниц. " \
                              f"Это может привести к поверхностному раскрытию темы каждого пункта."
                
            context.user_data["custom_plan"] = cleaned_plan
            context.user_data["plan_entered"] = True
            
            # Показываем введенный план пользователю
            plan_preview = "\n".join([f"{i+1}. {item}" for i, item in enumerate(cleaned_plan)])
            keyboard = create_keyboard_with_back([["⏭️ Пропустить"]], show_back=True)
            reply_markup = ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
            await update.message.reply_text(
                f"✅ План принят:\n\n{plan_preview}{warning_text}\n\n📝 Шаг 6/6: Опишите ваши предпочтения по работе (например, стиль, источники, сроки):",
                reply_markup=reply_markup
            )
            return PREFERENCES
            
        except ValueError as e:
            await update.message.reply_text(f"⚠️ Ошибка: {e}")
            return PREFERENCES
    else:
        # Обрабатываем предпочтения (для автоматического плана или финальный шаг для пользовательского)
        
        # Проверяем, нажал ли пользователь "Пропустить"
        if update.message.text == "⏭️ Пропустить":
            context.user_data["preferences"] = "Без особых предпочтений"
        else:
            try:
                preferences = validate_user_input(update.message.text, 500)
                context.user_data["preferences"] = preferences
            except ValueError as e:
                await update.message.reply_text(f"⚠️ Ошибка: {e}")
                return PREFERENCES
        
        # Переходим к оплате
        work_type = context.user_data["work_type"]
        
        # Извлекаем чистое название типа работы без эмодзи и цены
        clean_type = re.sub(r"[^А-Яа-яЁё ]", "", work_type).strip()
        
        # Определяем цену на основе чистого названия (всегда пересчитываем для надежности)
        if clean_type in ["Эссе", "Доклад"]:
            price = 300
        elif clean_type in ["Реферат", "Проект"]:
            price = 400
        elif clean_type == "Курсовая работа":
            price = 500
        elif clean_type == "Дипломная работа":
            price = 800
        else:
            price = 300
        context.user_data["price"] = price
        
        context.user_data["current_step"] = PAYMENT
        
        # Создание платежа через YooKassa (или тестовый заказ)
        keyboard = [
            [InlineKeyboardButton("💳 Оплатить заказ", callback_data="pay")],
            [InlineKeyboardButton("◀️ Назад", callback_data="back")]
        ]
        reply_markup = InlineKeyboardMarkup(keyboard)
        await update.message.reply_text(
            f"💰 *Оплата заказа* 💰\n\n"
            f"Сумма к оплате: {price} рублей\n\n"
            "Нажмите кнопку ниже для перехода к оплате:",
            reply_markup=reply_markup,
            parse_mode='Markdown'
        )
        return PAYMENT


async def create_payment(update: Update, context: CallbackContext) -> int:
    try:
        # Создаем заказ в базе данных
        user_id = update.callback_query.from_user.id
        username = update.callback_query.from_user.username or ""
        
        order_data = {
            "work_type": context.user_data.get("work_type", ""),
            "science_name": context.user_data.get("science_name", ""),
            "work_theme": context.user_data.get("work_theme", ""),
            "page_number": context.user_data.get("page_number", 0),
            "price": context.user_data.get("price", 0)
        }
        order_id = await create_order(user_id, order_data)
        context.user_data["order_id"] = order_id
        
        # Получаем данные из context
        price = context.user_data.get("price", 300)
        work_type = context.user_data.get("work_type", "Работа")
        
        # Проверяем, является ли пользователь тестовым
        is_test_user = username in TEST_MODE_USERNAMES
        
        # 🧪 ТЕСТОВЫЙ РЕЖИМ - пропускаем реальную оплату
        if TESTING_MODE or is_test_user:
            mode_reason = "тестовый пользователь @" + username if is_test_user else "глобальный тестовый режим"
            logging.info(f"🧪 ТЕСТОВЫЙ РЕЖИМ ({mode_reason}): Пропуск оплаты для заказа {order_id}")
            await update.callback_query.answer()
            await update.callback_query.edit_message_text(
                f"✅ Оплата успешно проведена! Начинаем выполнение вашего заказа.\n\n"
                f"🔄 Генерация вашей работы... Пожалуйста, подождите!\n"
                f"Это займет 2-3 минуты."
            )
            # Обновляем статус как оплаченный
            await update_order_status(order_id, "paid", payment_id="TEST_MODE")
            # Запускаем процесс генерации
            asyncio.create_task(process_order(context, user_id, order_id))
            return ConversationHandler.END
        
        # Проверяем настройки YooKassa для реального режима
        if not YOOKASSA_SHOP_ID or not YOOKASSA_SECRET_KEY:
            await update.callback_query.answer("❌ Платежная система временно недоступна")
            await update.callback_query.edit_message_text(
                "❌ Извините, платежная система временно недоступна. "
                "Пожалуйста, попробуйте позже или обратитесь в поддержку."
            )
            return ConversationHandler.END
        
        try:
            # Отладочное логирование перед созданием платежа
            logging.info(f"Создание платежа для заказа {order_id}, сумма: {price}")
            logging.info(f"Данные получателя: email=user{user_id}@ninjaessayai.com")
            
            payment = Payment.create({
                "amount": {"value": f"{price}.00", "currency": "RUB"},
                "confirmation": {"type": "redirect", "return_url": "https://t.me/NinjaEssayAI_bot"},
                "capture": True,
                "description": f"{work_type} - {context.user_data.get('work_theme', 'Тема не указана')}",
                "metadata": {"order_id": str(order_id)},
                "receipt": {
                    "customer": {
                        "email": f"user{user_id}@ninjaessayai.com"
                    },
                    "items": [
                        {
                            "description": f"Написание работы: {work_type}",
                            "quantity": "1.00",
                            "amount": {
                                "value": f"{price}.00",
                                "currency": "RUB"
                            },
                            "vat_code": 1,
                            "payment_mode": "full_payment",
                            "payment_subject": "service"
                        }
                    ]
                }
            }, uuid.uuid4())
            
            logging.info(f"Платеж создан успешно: {payment.id}")
            
            # Обновляем заказ с payment_id
            await update_order_status(order_id, "payment_created", payment.id)
            
            await update.callback_query.answer()
            
            # Создаем красивую кнопку оплаты
            payment_keyboard = InlineKeyboardMarkup([
                [InlineKeyboardButton("✅ Перейти к оплате", url=payment.confirmation.confirmation_url)]
            ])
            
            await update.callback_query.edit_message_text(
                f"💳 Оплата {price}₽\n\n"
                "После оплаты работа будет создана автоматически (2-3 минуты).",
                reply_markup=payment_keyboard
            )
            
            # Запускаем мониторинг платежа
            asyncio.create_task(monitor_payment(context, update.callback_query.from_user.id, payment.id))
            return ConversationHandler.END
            
        except Exception as payment_error:
            # Более подробное логирование ошибки
            error_details = str(payment_error)
            if hasattr(payment_error, '__dict__'):
                error_details = getattr(payment_error, '__dict__', str(payment_error))
            
            logging.error(f"Ошибка создания платежа: {payment_error}")
            logging.error(f"Детали ошибки: {error_details}")
            
            await update_order_status(order_id, "payment_failed")
            await update.callback_query.answer("❌ Ошибка создания платежа")
            await update.callback_query.edit_message_text(
                "❌ Произошла ошибка при создании платежа. "
                "Пожалуйста, попробуйте позже или обратитесь в поддержку."
            )
            return ConversationHandler.END
            
    except Exception as e:
        logging.error(f"Критическая ошибка в create_payment: {e}")
        try:
            await update.callback_query.answer("❌ Системная ошибка")
            await update.callback_query.edit_message_text(
                "❌ Произошла системная ошибка. Пожалуйста, обратитесь в поддержку."
            )
        except:
            pass
        return ConversationHandler.END

# Функция для обработки заказа в тестовом режиме
async def process_order(context: CallbackContext, chat_id: int, order_id: int) -> None:
    """Обработка заказа в тестовом режиме без реальной оплаты"""
    try:
        logging.info(f"🧪 ТЕСТОВЫЙ РЕЖИМ: Начало обработки заказа {order_id}")
        
        await context.bot.send_message(chat_id=chat_id, text="✅ Заказ принят! Начинаем выполнение вашего заказа.")
        await context.bot.send_message(chat_id=chat_id, text="🔄 Генерация вашей работы... Пожалуйста, подождите!")
        
        # Генерация плана
        plan_array = await generate_plan(context)
        if not plan_array:
            raise RuntimeError("План не сгенерирован")
        context.user_data["plan_array"] = plan_array
        
        # Генерация текста
        doc_io = await generate_text(plan_array, context)
        if doc_io is None:
            raise RuntimeError("Документ не создан")
        
        # Создаем безопасное имя файла
        work_type = context.user_data.get("work_type", "Работа")
        work_theme = context.user_data.get("work_theme", "Тема")
        safe_type = sanitize_filename(work_type)
        safe_theme = sanitize_filename(work_theme)
        filename = f"{safe_type}_{safe_theme}.docx"
        
        # Отправляем документ
        await context.bot.send_document(
            chat_id=chat_id, 
            document=doc_io, 
            filename=filename,
            caption="✅ Ваша работа готова!"
        )
        
        # Проверяем, является ли работа докладом - генерируем презентацию
        work_type_clean = re.sub(r"[^А-Яа-яЁё ]", "", work_type).strip()
        if "Доклад" in work_type_clean:
            try:
                await context.bot.send_message(
                    chat_id=chat_id,
                    text="🎨 Создаем презентацию PowerPoint... Это займет 1-2 минуты."
                )
                
                # Генерация презентации
                pptx_io = await generate_presentation(plan_array, context)
                pptx_filename = f"{safe_type}_{safe_theme}.pptx"
                
                # Отправляем презентацию
                await context.bot.send_document(
                    chat_id=chat_id,
                    document=pptx_io,
                    filename=pptx_filename,
                    caption="🎨 Презентация готова! Профессиональный дизайн с изображениями и графиками."
                )
                logging.info(f"✅ Презентация для заказа {order_id} создана")
            except Exception as pptx_error:
                logging.error(f"Ошибка создания презентации: {pptx_error}")
                await context.bot.send_message(
                    chat_id=chat_id,
                    text="⚠️ Презентация не создана из-за технической ошибки. Документ Word готов."
                )
        
        await context.bot.send_message(
            chat_id=chat_id, 
            text="🎉 Ваш заказ выполнен! Спасибо за использование нашего сервиса!"
        )
        
        # Обновляем статус заказа как выполненный
        await update_order_status(order_id, "completed")
        logging.info(f"🧪 ТЕСТОВЫЙ РЕЖИМ: Заказ {order_id} успешно завершён")
        
    except Exception as e:
        logging.error(f"🧪 ТЕСТОВЫЙ РЕЖИМ: Ошибка при обработке заказа {order_id}: {e}")
        
        # Обновляем статус заказа как неудачный
        await update_order_status(order_id, "failed")
        
        await context.bot.send_message(
            chat_id=chat_id, 
            text="❌ Произошла ошибка при генерации работы.\n"
                 "Попробуйте еще раз или обратитесь в поддержку."
        )

# Функция для автоматического мониторинга статуса платежа и начала выполнения заказа
async def monitor_payment(context: CallbackContext, chat_id: int, payment_id: str) -> None:
    try:
        order_id = context.user_data.get("order_id")
        
        while True:
            payment = Payment.find_one(payment_id)
            status = payment.status
            
            if status == "succeeded":
                # Обновляем статус заказа
                if order_id:
                    await update_order_status(order_id, "paid", payment_id)
                
                # Попытка генерации и возврат при ошибке
                try:
                    await context.bot.send_message(chat_id=chat_id, text="✅ Оплата успешно проведена! Начинаем выполнение вашего заказа.")
                    await context.bot.send_message(chat_id=chat_id, text="🔄 Генерация вашей работы... Пожалуйста, подождите!")
                    
                    plan_array = await generate_plan(context)
                    if not plan_array:
                        raise RuntimeError("План не сгенерирован")
                    context.user_data["plan_array"] = plan_array
                    
                    doc_io = await generate_text(plan_array, context)
                    if doc_io is None:
                        raise RuntimeError("Документ не создан")
                    
                    # Создаем безопасное имя файла
                    work_type = context.user_data.get("work_type", "Работа")
                    work_theme = context.user_data.get("work_theme", "Тема")
                    safe_type = sanitize_filename(work_type)
                    safe_theme = sanitize_filename(work_theme)
                    filename = f"{safe_type}_{safe_theme}.docx"
                    
                    await context.bot.send_document(
                        chat_id=chat_id, 
                        document=doc_io, 
                        filename=filename,
                        caption="✅ Ваша работа готова! Спасибо за использование NinjaEssayAI!"
                    )
                    
                    # Проверяем, является ли работа докладом - генерируем презентацию
                    work_type_clean = re.sub(r"[^А-Яа-яЁё ]", "", work_type).strip()
                    if "Доклад" in work_type_clean:
                        try:
                            await context.bot.send_message(
                                chat_id=chat_id,
                                text="🎨 Создаем презентацию PowerPoint... Это займет 1-2 минуты."
                            )
                            
                            # Генерация презентации
                            pptx_io = await generate_presentation(plan_array, context)
                            pptx_filename = f"{safe_type}_{safe_theme}.pptx"
                            
                            # Отправляем презентацию
                            await context.bot.send_document(
                                chat_id=chat_id,
                                document=pptx_io,
                                filename=pptx_filename,
                                caption="🎨 Презентация готова! Профессиональный дизайн с изображениями и графиками."
                            )
                            logging.info(f"✅ Презентация для заказа {order_id} создана")
                        except Exception as pptx_error:
                            logging.error(f"Ошибка создания презентации: {pptx_error}")
                            await context.bot.send_message(
                                chat_id=chat_id,
                                text="⚠️ Презентация не создана из-за технической ошибки. Документ Word готов."
                            )
                    
                    await context.bot.send_message(chat_id=chat_id, text="🎉 Ваш заказ выполнен! Спасибо за использование нашего сервиса!")
                    
                    # Обновляем статус заказа как выполненный
                    if order_id:
                        await update_order_status(order_id, "completed")
                    
                    return
                    
                except Exception as gen_error:
                    logging.error(f"Ошибка при генерации: {gen_error}")
                    
                    # Обновляем статус заказа как неудачный
                    if order_id:
                        await update_order_status(order_id, "failed")
                    
                    # Автоматический возврат средств
                    price = context.user_data.get("price")
                    if price is None and hasattr(payment, 'amount'):
                        try:
                            price = float(payment.amount.value)
                        except Exception:
                            price = None
                    amount_value = f"{float(price):.2f}" if price is not None else payment.amount.value
                    
                    try:
                        Refund.create({
                            "payment_id": payment_id,
                            "amount": {"value": amount_value, "currency": "RUB"}
                        }, uuid.uuid4())
                        
                        # Обновляем статус заказа как возвращенный
                        if order_id:
                            await update_order_status(order_id, "refunded")
                        
                        await context.bot.send_message(chat_id=chat_id, text="❌ Произошла ошибка при генерации. Платеж возвращён.")
                    except Exception as refund_error:
                        logging.error(f"Ошибка при возврате платежа: {refund_error}")
                        await context.bot.send_message(chat_id=chat_id, text="⚠️ Ошибка при возврате средств. Пожалуйста, обратитесь в поддержку.")
                    return
                    
            if status in ("canceled", "failed"):
                if order_id:
                    await update_order_status(order_id, "cancelled")
                await context.bot.send_message(chat_id=chat_id, text="❌ Платеж отменён или не прошёл. Заказ отменён.")
                return
                
            await asyncio.sleep(5)
            
    except Exception as e:
        logging.error(f"Ошибка при мониторинге платежа: {e}")
        await context.bot.send_message(chat_id=chat_id, text="❌ Произошла ошибка при проверке статуса платежа.")
        return

# Функция для обработки тестового заказа (без реальной оплаты)
async def create_test_order(update: Update, context: CallbackContext) -> int:
    """Создает тестовый заказ без реальной оплаты"""
    try:
        await update.callback_query.answer()
        await update.callback_query.edit_message_text(
            "🧪 Тестовый заказ создан!\n"
            "🔄 Начинаем генерацию вашей работы... Пожалуйста, подождите!"
        )
        
        # Создаем заказ в базе данных
        user_id = update.callback_query.from_user.id
        order_data = {
            "work_type": context.user_data.get("work_type", ""),
            "science_name": context.user_data.get("science_name", ""),
            "work_theme": context.user_data.get("work_theme", ""),
            "page_number": context.user_data.get("page_number", 0),
            "price": 0  # Бесплатно в тестовом режиме
        }
        order_id = await create_order(user_id, order_data)
        await update_order_status(order_id, "test_paid")  # Отмечаем как тестовый заказ
        
        # Сразу начинаем генерацию (без ожидания платежа)
        try:
            # Генерация плана
            plan_array = await generate_plan(context)
            if not plan_array:
                raise RuntimeError("План не сгенерирован")
            context.user_data["plan_array"] = plan_array
            
            # Генерация документа
            doc_io = await generate_text(plan_array, context)
            if doc_io is None:
                raise RuntimeError("Документ не создан")
            
            # Создаем безопасное имя файла
            work_type = context.user_data.get("work_type", "Работа")
            work_theme = context.user_data.get("work_theme", "Тема")
            safe_type = sanitize_filename(work_type)
            safe_theme = sanitize_filename(work_theme)
            filename = f"{safe_type}_{safe_theme}.docx"
            
            # Отправляем готовый документ
            await update.callback_query.message.reply_document(
                document=doc_io, 
                filename=filename,
                caption="✅ Ваша работа готова! (Тестовый режим)\n🎉 Спасибо за использование NinjaEssayAI!"
            )
            
            # Обновляем статус заказа как выполненный
            await update_order_status(order_id, "test_completed")
            
        except Exception as gen_error:
            logging.error(f"Ошибка при генерации в тестовом режиме: {gen_error}")
            await update_order_status(order_id, "test_failed")
            await update.callback_query.message.reply_text(
                "❌ Произошла ошибка при генерации работы.\n"
                "Попробуйте еще раз или обратитесь в поддержку."
            )
        
        return ConversationHandler.END
        
    except Exception as e:
        logging.error(f"Критическая ошибка в create_test_order: {e}")
        try:
            await update.callback_query.answer("❌ Системная ошибка")
            await update.callback_query.edit_message_text(
                "❌ Произошла системная ошибка. Пожалуйста, обратитесь в поддержку."
            )
        except:
            pass
        return ConversationHandler.END



# Обработка оплаты и генерация документа
# Упрощение интерактивного элемента загрузки
async def send_loading_message(update: Update, context: CallbackContext) -> None:
    await update.message.reply_text("🔄 Генерация вашей работы... Пожалуйста, подождите!")

# Вызов упрощенного элемента загрузки перед генерацией плана и текста
async def pay(update: Update, context: CallbackContext) -> int:
    logging.info("Начало выполнения заказа после оплаты.")
    await update.message.reply_text("Оплата успешно проведена! Начинаем выполнение вашего заказа.")

    # Упрощенный элемент загрузки
    await send_loading_message(update, context)

    # Генерация плана
    logging.info("Генерация плана работы.")
    plan_array = await generate_plan(context)
    context.user_data["plan_array"] = plan_array
    logging.info(f"План работы сгенерирован: {plan_array}")

    # Генерация текста и создание документа в памяти
    logging.info("Генерация текста по главам плана.")
    doc_io = await generate_text(plan_array, context)
    
    if doc_io is None:
        await update.message.reply_text("❌ Произошла ошибка при создании документа.")
        return ConversationHandler.END
    
    logging.info("Документ создан в памяти")

    # Создаем безопасное имя файла для отправки
    work_type = context.user_data.get("work_type", "Работа")
    work_theme = context.user_data.get("work_theme", "Тема")
    safe_type = sanitize_filename(work_type)
    safe_theme = sanitize_filename(work_theme)
    filename = f"{safe_type}_{safe_theme}.docx"

    # Отправка документа пользователю напрямую из памяти
    await update.message.reply_document(
        document=doc_io,
        filename=filename,
        caption="✅ Ваша работа готова! Спасибо за использование NinjaEssayAI!"
    )

    await update.message.reply_text("🎉 Заказ выполнен успешно! Спасибо за использование нашего сервиса!")
    logging.info("Заказ успешно выполнен и отправлен пользователю.")
    return ConversationHandler.END

# Enhanced error handling in the generate_plan function.
async def generate_plan(context: CallbackContext) -> list:
    # Проверяем, есть ли пользовательский план
    use_custom_plan = context.user_data.get("use_custom_plan", False)
    page_number = context.user_data.get("page_number", 0)
    
    if use_custom_plan:
        custom_plan = context.user_data.get("custom_plan", [])
        if custom_plan:
            # Проверяем соответствие количества пунктов плана заявленному количеству страниц
            expected_chapters = max(1, page_number // 2)  # Минимум 1 пункт
            plan_chapters = len(custom_plan)
            
            logging.info(f"Пользовательский план: {custom_plan}")
            logging.info(f"Заявлено страниц: {page_number}, пунктов в плане: {plan_chapters}, ожидалось: {expected_chapters}")
            
            # Если пунктов слишком мало для заявленного количества страниц, предупреждаем
            if plan_chapters < expected_chapters:
                try:
                    chat_id = get_chat_id(context)
                    await context.bot.send_message(
                        chat_id=chat_id,
                        text=f"⚠️ Внимание: Для {page_number} страниц рекомендуется {expected_chapters} пунктов плана, "
                             f"а у вас {plan_chapters}. Это может повлиять на объем работы."
                    )
                except Exception as chat_error:
                    logging.error(f"Ошибка отправки предупреждения: {chat_error}")
            
            return custom_plan
    
    # Если пользовательского плана нет, генерируем автоматически
    logging.info("Запрос на генерация плана отправлен в DeepSeek API.")
    
    science_name = context.user_data.get("science_name", "")
    work_type = context.user_data.get("work_type", "")
    work_theme = context.user_data.get("work_theme", "")
    preferences = context.user_data.get("preferences", "")
    page_number = context.user_data.get("page_number", 0)

    if not science_name or not work_type or not work_theme or not page_number:
        logging.error("Недостаточно данных для генерации плана.")
        try:
            chat_id = get_chat_id(context)
            await context.bot.send_message(
                chat_id=chat_id,
                text="Ошибка: Недостаточно данных для генерации плана. Проверьте введенные данные."
            )
        except Exception as chat_error:
            logging.error(f"Ошибка отправки сообщения: {chat_error}")
        return []

    # Минимум 3 пункта (Введение + 1 глава + Заключение)
    # Расчет: страницы / 2, но не менее 3
    calls_number = max(3, page_number // 2)

    prompt = (
        f"Действуй как специалист в области {science_name}. "
        f"Составь подробный план из {calls_number} пунктов для {work_type} "
        f"по теме: {work_theme}. Учти предпочтения: {preferences}. "
        "ОБЯЗАТЕЛЬНО включи в план:\n"
        "1. 'Введение' (первый пункт)\n"
        f"2-{calls_number-1}. Основные разделы (содержательные названия глав)\n"
        f"{calls_number}. 'Заключение' (последний пункт)\n\n"
        "Верни ТОЛЬКО нумерованный список без дополнительных комментариев:\n"
        "1. Введение\n"
        "2. [Название раздела 1]\n"
        "3. [Название раздела 2]\n"
        f"{calls_number}. Заключение"
    )

    try:
        response = await client.chat.completions.create(
            model="deepseek-reasoner",
            messages=[
                {"role": "system", "content": "mode: plan_generation"},
                {"role": "user", "content": prompt}
            ],
            stream=False
        )
        response_content = response.choices[0].message.content
    except Exception as e:
        logging.error(f"Ошибка при вызове DeepSeek API: {e}")
        try:
            chat_id = get_chat_id(context)
            await context.bot.send_message(
                chat_id=chat_id,
                text="Произошла ошибка при генерации плана. Пожалуйста, попробуйте позже."
            )
        except Exception as chat_error:
            logging.error(f"Ошибка отправки сообщения: {chat_error}")
        return []

    logging.info("Ответ от DeepSeek API получен.")

    try:
        plan_array = json.loads(response_content)
        if not isinstance(plan_array, list):
            raise ValueError("Ответ не является массивом.")
    except (json.JSONDecodeError, ValueError):
        logging.info("Ответ не в формате JSON, обрабатываем как текст.")
        lines = response_content.splitlines()
        plan_array = [line.strip() for line in lines if line.strip()]
        # Удаляем нумерацию и очищаем
        plan_array = [re.sub(r'^\d+\.\s*', '', item).strip() for item in plan_array if item.strip()]
        # Фильтруем пустые строки и мусор
        plan_array = [item for item in plan_array if item and len(item) > 2]

    # Проверяем и дополняем план
    if len(plan_array) < calls_number:
        # Если план слишком короткий, дополняем
        needed = calls_number - len(plan_array)
        for i in range(needed):
            plan_array.insert(len(plan_array), f"Раздел {len(plan_array)}")
    elif len(plan_array) > calls_number:
        # Если план слишком длинный, обрезаем
        plan_array = plan_array[:calls_number]
    
    # Удаляем эмодзи из всех элементов плана
    plan_array = [remove_emojis(item) for item in plan_array]
    
    # Гарантируем наличие Введения и Заключения
    if plan_array and not any(keyword in plan_array[0].lower() for keyword in ['введение', 'introduction']):
        plan_array[0] = 'Введение'
    
    if plan_array and not any(keyword in plan_array[-1].lower() for keyword in ['заключение', 'conclusion', 'выводы']):
        plan_array[-1] = 'Заключение'
    
    # Если план пустой или слишком короткий, создаем минимальный план
    if not plan_array or len(plan_array) < 3:
        logging.warning("План от DeepSeek некорректен, создаем базовый план")
        plan_array = ['Введение']
        for i in range(calls_number - 2):
            plan_array.append(f'Глава {i+1}')
        plan_array.append('Заключение')

    logging.info(f"Итоговый план: {plan_array}")
    return plan_array

# Функция для добавления нумерации страниц
def add_page_number(section, start_number=1):
    """Добавляет нумерацию страниц в раздел документа
    
    Args:
        section: Раздел документа
        start_number: Начальный номер страницы (по умолчанию 1)
    """
    footer = section.footer
    footer_paragraph = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
    footer_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = footer_paragraph.add_run()
    
    # Устанавливаем начальный номер страницы для раздела
    if start_number > 1:
        sectPr = section._sectPr
        pgNumType = sectPr.find(qn('w:pgNumType'))
        if pgNumType is None:
            pgNumType = OxmlElement('w:pgNumType')
            sectPr.insert(0, pgNumType)
        pgNumType.set(qn('w:start'), str(start_number))
    
    fldChar1 = OxmlElement('w:fldChar')
    fldChar1.set(qn('w:fldCharType'), 'begin')
    run._r.append(fldChar1)
    instrText = OxmlElement('w:instrText')
    instrText.text = 'PAGE'
    run._r.append(instrText)
    fldChar2 = OxmlElement('w:fldChar')
    fldChar2.set(qn('w:fldCharType'), 'end')
    run._r.append(fldChar2)

# Функция для добавления титульного листа
def add_title_page(doc, work_type, work_theme, science_name, page_number):
    """Добавляет титульный лист в документ"""
    
    # Очищаем тему от смайликов
    import re
    emoji_pattern = re.compile(
        "["
        u"\U0001F600-\U0001F64F"
        u"\U0001F300-\U0001F5FF"
        u"\U0001F680-\U0001F6FF"
        u"\U0001F1E0-\U0001F1FF"
        u"\U00002702-\U000027B0"
        u"\U000024C2-\U0001F251"
        u"\U0001F900-\U0001F9FF"
        u"\U0001FA70-\U0001FAFF"
        "]+", flags=re.UNICODE
    )
    work_theme = emoji_pattern.sub('', work_theme).strip()
    
    # Название учебного заведения
    university_p = doc.add_paragraph()
    university_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    university_run = university_p.add_run("МИНИСТЕРСТВО ОБРАЗОВАНИЯ И НАУКИ РОССИЙСКОЙ ФЕДЕРАЦИИ\n")
    university_run.font.name = 'Times New Roman'
    university_run.font.size = Pt(14)
    university_run.font.bold = True
    
    university_run2 = university_p.add_run("ФЕДЕРАЛЬНОЕ ГОСУДАРСТВЕННОЕ БЮДЖЕТНОЕ ОБРАЗОВАТЕЛЬНОЕ УЧРЕЖДЕНИЕ\n")
    university_run2.font.name = 'Times New Roman'
    university_run2.font.size = Pt(14)
    university_run2.font.bold = True
    
    university_run3 = university_p.add_run("ВЫСШЕГО ОБРАЗОВАНИЯ\n")
    university_run3.font.name = 'Times New Roman'
    university_run3.font.size = Pt(14)
    university_run3.font.bold = True
    
    university_run4 = university_p.add_run("«РОССИЙСКИЙ УНИВЕРСИТЕТ»")
    university_run4.font.name = 'Times New Roman'
    university_run4.font.size = Pt(14)
    university_run4.font.bold = True
    
    # Добавляем отступ
    doc.add_paragraph()
    doc.add_paragraph()
    doc.add_paragraph()
    
    # Кафедра
    department_p = doc.add_paragraph()
    department_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    department_run = department_p.add_run(f"Кафедра {science_name}")
    department_run.font.name = 'Times New Roman'
    department_run.font.size = Pt(14)
    
    # Добавляем отступ
    doc.add_paragraph()
    doc.add_paragraph()
    
    # Тип работы
    work_type_p = doc.add_paragraph()
    work_type_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    work_type_run = work_type_p.add_run(work_type.upper())
    work_type_run.font.name = 'Times New Roman'
    work_type_run.font.size = Pt(16)
    work_type_run.font.bold = True
    
    # Добавляем отступ
    doc.add_paragraph()
    
    # Тема работы
    theme_p = doc.add_paragraph()
    theme_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    theme_run = theme_p.add_run(f"на тему: «{work_theme}»")
    theme_run.font.name = 'Times New Roman'
    theme_run.font.size = Pt(14)
    theme_run.font.bold = True
    
    # Добавляем отступ
    doc.add_paragraph()
    doc.add_paragraph()
    doc.add_paragraph()
    doc.add_paragraph()
    
    # Информация о студенте и преподавателе (справа)
    info_table = doc.add_table(rows=6, cols=2)
    info_table.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    
    # Настраиваем таблицу
    for row in info_table.rows:
        for cell in row.cells:
            for paragraph in cell.paragraphs:
                paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
    
    # Заполняем таблицу
    info_table.cell(0, 0).text = "Дисциплина:"
    info_table.cell(0, 1).text = science_name
    
    info_table.cell(1, 0).text = "Выполнил(а):"
    info_table.cell(1, 1).text = "студент(ка) группы ___________"
    
    info_table.cell(2, 0).text = ""
    info_table.cell(2, 1).text = "_________________________"
    
    info_table.cell(3, 0).text = "Проверил:"
    info_table.cell(3, 1).text = "_________________________"
    
    info_table.cell(4, 0).text = ""
    info_table.cell(4, 1).text = "(должность, ученая степень, звание)"
    
    info_table.cell(5, 0).text = ""
    info_table.cell(5, 1).text = "_________________________"
    
    # Настраиваем шрифт в таблице
    for row in info_table.rows:
        for cell in row.cells:
            for paragraph in cell.paragraphs:
                for run in paragraph.runs:
                    run.font.name = 'Times New Roman'
                    run.font.size = Pt(12)
    
    # Добавляем отступ
    doc.add_paragraph()
    doc.add_paragraph()
    doc.add_paragraph()
    
    # Год и город
    footer_p = doc.add_paragraph()
    footer_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer_run = footer_p.add_run("Москва 2025")
    footer_run.font.name = 'Times New Roman'
    footer_run.font.size = Pt(14)


def parse_preferences_by_chapter(preferences: str, plan_array: list) -> dict:
    """
    Парсит пожелания пользователя и распределяет их по главам
    
    Работает в несколько этапов:
    1. Парсинг по явным ключевым словам (введение, заключение, основная часть)
    2. Интеллектуальное распределение по контентным ключевым словам (цель, задачи, методология и т.д.)
    3. Сопоставление с фактическими названиями глав из плана
    
    Args:
        preferences: Текст пожеланий пользователя
        plan_array: Массив глав плана (может быть автоматическим или пользовательским)
    
    Returns:
        Словарь: {
            'global': 'общие пожелания для всех глав',
            'by_chapter': {
                'Введение': 'специфичные пожелания для введения',
                'Глава 1': 'специфичные пожелания для главы 1',
                ...
            }
        }
    """
    result = {
        'global': '',
        'by_chapter': {}
    }
    
    if not preferences or preferences == "Без особых предпочтений":
        return result
    
    # === ЭТАП 1: Базовые ключевые слова для типовых разделов ===
    base_chapter_keywords = {
        'введение': ['введение', 'введении', 'во введение', 'для введения', 'во введении'],
        'заключение': ['заключение', 'заключении', 'в заключение', 'для заключения', 'в заключении'],
        'основная часть': ['основная часть', 'основной части', 'в основной части', 'основную часть'],
    }
    
    # === ЭТАП 2: Контентные ключевые слова (по смыслу содержания) ===
    content_keywords = {
        'введение': [
            'цель', 'задачи', 'задач', 'актуальность', 'актуальности', 
            'проблема', 'проблемы', 'постановка задачи',
            'объект исследования', 'предмет исследования', 'гипотеза', 'гипотезу'
        ],
        'основная часть': [
            'методология', 'методологию', 'анализ', 'исследование', 'практика', 
            'теория', 'теоретическая часть', 'обзор литературы', 
            'эксперимент', 'данные', 'данных', 'результаты', 'результатов'
        ],
        'заключение': [
            'выводы', 'вывод', 'итоги', 'итог', 'резюме', 'достижения', 
            'перспективы', 'перспектив', 'рекомендации', 'рекомендаций',
            'заключительные положения'
        ]
    }
    
    # === ЭТАП 3: Создание словаря ключевых слов для каждой главы плана ===
    all_chapter_keywords = {}
    for chapter in plan_array:
        chapter_lower = chapter.lower()
        all_chapter_keywords[chapter_lower] = [chapter_lower]
        
        # Добавляем базовые ключевые слова
        for base_chapter, keywords in base_chapter_keywords.items():
            if base_chapter in chapter_lower:
                all_chapter_keywords[chapter_lower].extend(keywords)
                # Добавляем контентные ключевые слова
                if base_chapter in content_keywords:
                    all_chapter_keywords[chapter_lower].extend(content_keywords[base_chapter])
    
    # === ЭТАП 4: Разбиваем пожелания на предложения ===
    sentences = re.split(r'[.;!?]', preferences)
    
    for sentence in sentences:
        sentence = sentence.strip()
        if not sentence or len(sentence) < 5:
            continue
        
        sentence_lower = sentence.lower()
        matched = False
        
        # === ЭТАП 5: Поиск соответствий по ключевым словам ===
        for chapter in plan_array:
            chapter_lower = chapter.lower()
            keywords = all_chapter_keywords.get(chapter_lower, [])
            
            # Проверяем наличие любого ключевого слова в предложении
            for keyword in keywords:
                if keyword in sentence_lower:
                    if chapter not in result['by_chapter']:
                        result['by_chapter'][chapter] = []
                    
                    # Убираем упоминание главы из текста пожелания
                    cleaned = re.sub(
                        f'({"|".join(re.escape(k) for k in keywords)})',
                        '',
                        sentence,
                        flags=re.IGNORECASE
                    ).strip()
                    
                    # Очищаем от лишних предлогов в начале
                    cleaned = re.sub(r'^(в|для|по|о|про)\s+', '', cleaned, flags=re.IGNORECASE).strip()
                    
                    if cleaned and len(cleaned) > 5:
                        result['by_chapter'][chapter].append(cleaned)
                    matched = True
                    break
            
            if matched:
                break
        
        # Если не нашли соответствий - это общее пожелание для всей работы
        if not matched:
            if result['global']:
                result['global'] += '. ' + sentence
            else:
                result['global'] = sentence
    
    # === ЭТАП 6: Объединяем списки пожеланий для каждой главы в строку ===
    for chapter in result['by_chapter']:
        if isinstance(result['by_chapter'][chapter], list):
            result['by_chapter'][chapter] = '. '.join(result['by_chapter'][chapter])
    
    # === ЛОГИРОВАНИЕ ДЛЯ ОТЛАДКИ ===
    logging.info(f"📋 Умное распределение пожеланий по главам:")
    logging.info(f"   Общие пожелания: {result['global']}")
    for chapter, prefs in result['by_chapter'].items():
        logging.info(f"   {chapter}: {prefs}")
    
    return result


# ===================== ФУНКЦИИ ДЛЯ ГЕНЕРАЦИИ ПРЕЗЕНТАЦИЙ =====================

async def download_image_from_url(url: str) -> io.BytesIO:
    """Скачивает изображение по URL и возвращает его в памяти"""
    try:
        async with aiohttp.ClientSession() as session:
            async with session.get(url, timeout=aiohttp.ClientTimeout(total=10)) as response:
                if response.status == 200:
                    image_data = await response.read()
                    return io.BytesIO(image_data)
    except Exception as e:
        logging.error(f"Ошибка загрузки изображения {url}: {e}")
    return None

async def search_image_for_topic(topic: str) -> str:
    """Генерирует URL изображения с Unsplash API для темы"""
    # Используем Unsplash Source API для получения случайного изображения по теме
    # Это бесплатный сервис, не требующий API ключа
    query = topic.replace(' ', '+')
    url = f"https://source.unsplash.com/800x600/?{query}"
    return url

def create_chart_image(data: dict, chart_type: str = 'bar') -> io.BytesIO:
    """Создает изображение графика с помощью matplotlib"""
    try:
        fig, ax = plt.subplots(figsize=(8, 6))
        
        if chart_type == 'bar':
            ax.bar(data.keys(), data.values(), color='#4472C4')
        elif chart_type == 'line':
            ax.plot(list(data.keys()), list(data.values()), marker='o', color='#4472C4')
        elif chart_type == 'pie':
            ax.pie(data.values(), labels=data.keys(), autopct='%1.1f%%', startangle=90)
        
        ax.set_title('Визуализация данных', fontsize=14, fontweight='bold')
        plt.tight_layout()
        
        # Сохраняем в память
        img_buffer = io.BytesIO()
        plt.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight')
        img_buffer.seek(0)
        plt.close(fig)
        
        return img_buffer
    except Exception as e:
        logging.error(f"Ошибка создания графика: {e}")
        return None

async def extract_presentation_content(text_content: str, chapter_title: str) -> dict:
    """Извлекает ключевые пункты из текста для презентации с помощью DeepSeek"""
    try:
        prompt = f"""Из следующего текста выдели 3-5 ключевых пунктов для слайда презентации.
        
Текст главы "{chapter_title}":
{text_content[:1500]}

Верни ТОЛЬКО JSON в формате:
{{
    "title": "Заголовок слайда",
    "bullets": ["Пункт 1", "Пункт 2", "Пункт 3"],
    "chart_data": {{"Категория1": 30, "Категория2": 45, "Категория3": 25}}
}}

Если данных для графика нет, используй chart_data: null."""

        response = await client.chat.completions.create(
            model="deepseek-chat",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7,
            max_tokens=500
        )
        
        content = response.choices[0].message.content.strip()
        # Убираем markdown форматирование если есть
        if content.startswith('```'):
            content = content.split('\n', 1)[1]
            content = content.rsplit('```', 1)[0]
        
        return json.loads(content)
    except Exception as e:
        logging.error(f"Ошибка извлечения контента: {e}")
        return {
            "title": chapter_title,
            "bullets": ["Основные моменты главы", "Ключевые выводы", "Важные аспекты"],
            "chart_data": None
        }

def apply_modern_theme(prs: Presentation):
    """Применяет современный дизайн к презентации"""
    # Устанавливаем размер слайда (16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

def add_title_slide(prs: Presentation, title: str, subtitle: str, author: str = "NinjaEssayAI"):
    """Добавляет титульный слайд"""
    slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Заголовок
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = PptxPt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = PptxRGBColor(0, 0, 0)
    
    # Подзаголовок
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = f"{subtitle}\n\n{author}"
    subtitle_shape.text_frame.paragraphs[0].font.size = PptxPt(24)
    
    return slide

def add_content_slide(prs: Presentation, title: str, bullets: list, image_stream: io.BytesIO = None):
    """Добавляет слайд с контентом и изображением"""
    # Используем layout с заголовком и содержимым
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    # Заголовок
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = PptxPt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = PptxRGBColor(68, 114, 196)
    
    # Если есть изображение, используем layout с двумя колонками
    if image_stream:
        # Текст слева
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(4.5)
        height = Inches(3.5)
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        for bullet_text in bullets:
            p = text_frame.add_paragraph()
            p.text = bullet_text
            p.level = 0
            p.font.size = PptxPt(16)
            p.space_before = PptxPt(12)
        
        # Изображение справа
        image_stream.seek(0)
        left = Inches(5.5)
        top = Inches(1.5)
        slide.shapes.add_picture(image_stream, left, top, width=Inches(4), height=Inches(3.5))
    else:
        # Только текст
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        for bullet_text in bullets:
            p = text_frame.add_paragraph()
            p.text = bullet_text
            p.level = 0
            p.font.size = PptxPt(18)
            p.space_before = PptxPt(12)
    
    return slide

def add_chart_slide(prs: Presentation, title: str, chart_data: dict):
    """Добавляет слайд с графиком"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Заголовок
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = PptxPt(32)
    p.font.bold = True
    p.font.color.rgb = PptxRGBColor(68, 114, 196)
    
    # График
    chart_img = create_chart_image(chart_data, 'bar')
    if chart_img:
        left = Inches(1.5)
        top = Inches(1.5)
        slide.shapes.add_picture(chart_img, left, top, width=Inches(7), height=Inches(3.5))
    
    return slide

def add_conclusion_slide(prs: Presentation, title: str = "Спасибо за внимание!"):
    """Добавляет финальный слайд"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Центральный текст
    left = Inches(2)
    top = Inches(2)
    width = Inches(6)
    height = Inches(1.5)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = title
    
    p = text_frame.paragraphs[0]
    p.font.size = PptxPt(48)
    p.font.bold = True
    p.font.color.rgb = PptxRGBColor(68, 114, 196)
    p.alignment = PP_ALIGN.CENTER
    
    # Подпись
    left = Inches(3)
    top = Inches(4)
    width = Inches(4)
    height = Inches(0.5)
    
    footer_box = slide.shapes.add_textbox(left, top, width, height)
    footer_frame = footer_box.text_frame
    footer_frame.text = "Создано с помощью NinjaEssayAI"
    
    p = footer_frame.paragraphs[0]
    p.font.size = PptxPt(14)
    p.font.color.rgb = PptxRGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

async def generate_presentation(plan_array: list, context: CallbackContext) -> io.BytesIO:
    """Генерирует презентацию PowerPoint на основе плана работы"""
    try:
        logging.info("🎨 Начало генерации презентации")
        
        # Создаем презентацию
        prs = Presentation()
        apply_modern_theme(prs)
        
        # Получаем данные из контекста
        work_type = context.user_data.get("work_type", "Презентация")
        work_theme = context.user_data.get("work_theme", "Тема работы")
        science_name = context.user_data.get("science_name", "Дисциплина")
        
        # Титульный слайд
        add_title_slide(prs, work_theme, f"{work_type} по дисциплине: {science_name}")
        
        # Генерируем текст для каждой главы и создаем слайды
        for i, chapter_title in enumerate(plan_array):
            logging.info(f"📊 Создание слайда для главы: {chapter_title}")
            
            # Генерируем текст главы (упрощенная версия)
            chapter_text = await generate_chapter_text(chapter_title, context, i, len(plan_array))
            
            # Извлекаем ключевые пункты для слайда
            slide_content = await extract_presentation_content(chapter_text, chapter_title)
            
            # Пытаемся найти изображение для темы
            image_stream = None
            try:
                image_url = await search_image_for_topic(chapter_title)
                image_stream = await download_image_from_url(image_url)
            except Exception as e:
                logging.warning(f"Не удалось загрузить изображение: {e}")
            
            # Добавляем слайд с контентом
            add_content_slide(
                prs,
                slide_content['title'],
                slide_content['bullets'],
                image_stream
            )
            
            # Если есть данные для графика, добавляем слайд с графиком
            if slide_content.get('chart_data'):
                add_chart_slide(
                    prs,
                    f"{slide_content['title']} - Визуализация",
                    slide_content['chart_data']
                )
        
        # Заключительный слайд
        add_conclusion_slide(prs)
        
        # Сохраняем в память
        pptx_io = io.BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)
        
        logging.info("✅ Презентация успешно создана")
        return pptx_io
        
    except Exception as e:
        logging.error(f"❌ Ошибка генерации презентации: {e}")
        raise

async def generate_chapter_text(chapter_title: str, context: CallbackContext, chapter_index: int, total_chapters: int) -> str:
    """Генерирует текст для главы (упрощенная версия для презентации)"""
    try:
        science_name = context.user_data.get("science_name", "")
        work_theme = context.user_data.get("work_theme", "")
        preferences = context.user_data.get("preferences", "")
        
        # Определяем тип главы
        if chapter_index == 0 or "введение" in chapter_title.lower():
            chapter_type = "введение"
        elif chapter_index == total_chapters - 1 or "заключение" in chapter_title.lower():
            chapter_type = "заключение"
        else:
            chapter_type = "основная часть"
        
        prompt = f"""Напиши краткий текст (200-300 слов) для главы "{chapter_title}" 
по теме "{work_theme}" в области {science_name}.

Тип главы: {chapter_type}
Предпочтения: {preferences}

Сделай текст информативным и структурированным."""
        
        async with GENERATION_SEMAPHORE:
            response = await client.chat.completions.create(
                model="deepseek-reasoner",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.7,
                max_tokens=500
            )
        
        content = response.choices[0].message.content
        # Извлекаем текст из reasoning_content если есть
        if hasattr(response.choices[0].message, 'reasoning_content'):
            content = response.choices[0].message.reasoning_content or content
        
        return content.strip()
        
    except Exception as e:
        logging.error(f"Ошибка генерации текста главы: {e}")
        return f"Основное содержание главы {chapter_title}"


# Генерация текста и создание документа (в памяти)
async def generate_text(plan_array, context: CallbackContext) -> io.BytesIO:
    logging.info("Начало генерации текста по главам плана.")
    science_name = context.user_data["science_name"]
    work_type = context.user_data["work_type"]
    work_theme = context.user_data["work_theme"]
    preferences = context.user_data["preferences"]
    page_number = context.user_data.get("page_number", 0)

    if not plan_array:
        logging.error("План пуст. Невозможно сгенерировать текст.")
        try:
            chat_id = get_chat_id(context)
            await context.bot.send_message(
                chat_id=chat_id,
                text="Ошибка: План пуст. Убедитесь, что DeepSeek вернул корректный массив."
            )
        except Exception as chat_error:
            logging.error(f"Ошибка отправки сообщения: {chat_error}")
        return None

    # Рассчитываем количество слов на главу на основе заявленного количества страниц
    # Примерно 250-300 слов на страницу, распределяем равномерно между главами
    total_words = page_number * 275  # Среднее количество слов на страницу
    words_per_chapter = total_words // len(plan_array)  # Распределяем слова между главами
    
    # Для очень маленьких работ (1-2 страницы) минимум 200 слов на главу
    # Для средних работ (3-5 страниц) минимум 300 слов
    # Для больших работ (6+ страниц) минимум 400 слов
    if page_number <= 2:
        words_per_chapter = max(200, words_per_chapter)
    elif page_number <= 5:
        words_per_chapter = max(300, words_per_chapter)
    else:
        words_per_chapter = max(400, words_per_chapter)
    
    logging.info(f"Запрошено страниц: {page_number}, глав в плане: {len(plan_array)}, слов на главу: {words_per_chapter}")

    # === УМНОЕ РАСПРЕДЕЛЕНИЕ ПОЖЕЛАНИЙ ПО ГЛАВАМ ===
    parsed_preferences = parse_preferences_by_chapter(preferences, plan_array)
    global_preferences = parsed_preferences['global']
    chapter_preferences = parsed_preferences['by_chapter']

    # Функция для выполнения одного запроса к API
    async def fetch_chapter_text(chapter: str) -> tuple[str, str]:
        logging.info(f"Генерация текста для главы: {chapter}")
        
        # Определяем пожелания для этой конкретной главы
        chapter_specific = chapter_preferences.get(chapter, '')
        
        # Комбинируем общие и специфичные пожелания
        if chapter_specific and global_preferences:
            combined_preferences = f"{global_preferences}. {chapter_specific}"
        elif chapter_specific:
            combined_preferences = chapter_specific
        elif global_preferences:
            combined_preferences = global_preferences
        else:
            combined_preferences = preferences  # Используем оригинальные, если парсинг не дал результата
        
        prompt = (
            f"Действуй как специалист в области {science_name}, "
            f"напиши, строго с опорой на авторитетные источники, "
            f"главу: {chapter} в контексте написания {work_type} "
            f"по теме: {work_theme} (напиши не менее {words_per_chapter} слов) "
            f"(Напиши текст, в котором предложения будут иметь разную длину, "
            f"а также будет избегаться нахождение однокоренных слов "
            f"в соседних предложениях) "
            f"(избегай комментариев, анализа соблюденных тобою требований, "
            f"возвращай исключительно текст, так, будто бы ты отправляешь "
            f"его на проверку преподавателю) "
            f"НЕ используй вводные фразы типа 'Отлично', 'Вот текст', 'Рассмотрим'. "
            f"Начинай сразу с основного содержания. {combined_preferences}"
        )
        # Выполняем запрос к DeepSeek под контролем семафора и обрабатываем ошибки
        try:
            async with GENERATION_SEMAPHORE:
                response = await client.chat.completions.create(
                    model="deepseek-reasoner",
                    messages=[{"role": "user", "content": prompt}],
                    stream=False
                )
            chapter_text = response.choices[0].message.content
            # Валидируем и очищаем сгенерированный контент
            chapter_text = validate_generated_content(chapter_text, chapter)
            logging.info(f"Сгенерирован текст для главы: {chapter_text[:100]}...")
            return chapter, chapter_text
        except Exception as e:
            logging.error(f"Ошибка при генерации текста для главы {chapter}: {e}")
            # Возвращаем сообщение об ошибке, чтобы не сломать gather
            return chapter, f"Ошибка при генерации текста для главы: {chapter}"

    # Создание списка задач для параллельного выполнения
    tasks = [fetch_chapter_text(chapter) for chapter in plan_array]
    # Параллельное выполнение запросов
    results = await asyncio.gather(*tasks, return_exceptions=True)

    # Проверка результатов
    chapters_text = []
    for result in results:
        if isinstance(result, Exception):
            logging.error(f"Ошибка в задаче: {result}")
            try:
                chat_id = get_chat_id(context)
                await context.bot.send_message(
                    chat_id=chat_id,
                    text=f"Ошибка при генерации текста. Пожалуйста, попробуйте позже."
                )
            except Exception as chat_error:
                logging.error(f"Ошибка отправки сообщения: {chat_error}")
            return None
        elif isinstance(result, tuple) and len(result) == 2:
            chapter, chapter_text = result
            chapters_text.append((chapter, chapter_text))
        else:
            logging.error(f"Неожиданный формат результата: {result}")
            return None

    # Создание документа в памяти
    doc = docx.Document()
    
    # === ТИТУЛЬНЫЙ ЛИСТ (без нумерации) ===
    title_section = doc.sections[0]
    # Настройка полей по требованиям: верх/низ 25мм, лево 30мм, право 10мм
    title_section.top_margin = Cm(2.5)     # 25 мм
    title_section.bottom_margin = Cm(2.5)  # 25 мм
    title_section.left_margin = Cm(3.0)    # 30 мм
    title_section.right_margin = Cm(1.0)   # 10 мм
    
    # Отключаем нумерацию на титульном листе
    title_section.different_first_page_header_footer = True

    # Установка стиля текста для всего документа
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Times New Roman'
    font.size = Pt(14)
    
    # Настройка параграфа для стиля Normal
    paragraph_format = style.paragraph_format
    paragraph_format.line_spacing = 1.5  # Интервал 1.5
    paragraph_format.first_line_indent = Cm(1.25)  # Отступ первой строки

    # Добавляем титульный лист (без нумерации)
    add_title_page(doc, work_type, work_theme, science_name, context.user_data.get("page_number", 0))
    
    # Добавляем разрыв раздела после титульного листа
    doc.add_section()
    
    # === ОГЛАВЛЕНИЕ (начинается нумерация с 2) ===
    content_section = doc.sections[-1]
    content_section.top_margin = Cm(2.5)
    content_section.bottom_margin = Cm(2.5)
    content_section.left_margin = Cm(3.0)
    content_section.right_margin = Cm(1.0)
    
    # Добавляем нумерацию страниц со 2-й страницы
    add_page_number(content_section, start_number=2)
    
    # Добавляем заголовок "Оглавление"
    contents_heading = doc.add_heading("ОГЛАВЛЕНИЕ", level=1)
    contents_heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in contents_heading.runs:
        run.font.name = 'Times New Roman'
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0, 0, 0)  # Черный цвет
    
    # Добавляем пункты оглавления
    page_counter = 3  # Первая страница после оглавления
    
    # Главы основной части
    for i, chapter in enumerate(plan_array, 1):
        contents_p = doc.add_paragraph()
        contents_p.paragraph_format.first_line_indent = Cm(0)
        contents_p.paragraph_format.left_indent = Cm(0)
        
        # Добавляем номер и название главы
        chapter_text = f"{i}. {chapter}"
        run = contents_p.add_run(chapter_text)
        run.font.name = 'Times New Roman'
        run.font.size = Pt(14)
        
        # Добавляем точки-заполнители
        dots_count = max(1, 70 - len(chapter_text))
        dots_run = contents_p.add_run("." * dots_count)
        dots_run.font.name = 'Times New Roman'
        dots_run.font.size = Pt(14)
        
        # Номер страницы
        page_run = contents_p.add_run(f" {page_counter}")
        page_run.font.name = 'Times New Roman'
        page_run.font.size = Pt(14)
        page_counter += 1
    
    # Список источников
    contents_p = doc.add_paragraph()
    contents_p.paragraph_format.first_line_indent = Cm(0)
    contents_p.paragraph_format.left_indent = Cm(0)
    run = contents_p.add_run("Список источников")
    run.font.name = 'Times New Roman'
    run.font.size = Pt(14)
    dots_run = contents_p.add_run("." * 57)
    dots_run.font.name = 'Times New Roman'
    dots_run.font.size = Pt(14)
    page_run = contents_p.add_run(f" {page_counter}")
    page_run.font.name = 'Times New Roman'
    page_run.font.size = Pt(14)
    
    # Добавляем разрыв страницы после оглавления
    doc.add_page_break()

    # === ОСНОВНАЯ ЧАСТЬ ===
    # Добавление текста глав в документ
    for i, (chapter, chapter_text) in enumerate(chapters_text, 1):
        # Удаляем дублирующийся заголовок из текста главы
        chapter_text_cleaned = remove_chapter_title_from_text(chapter_text, chapter)
        
        chapter_heading = doc.add_heading(f"{i}. {chapter}", level=2)
        chapter_heading.alignment = WD_ALIGN_PARAGRAPH.LEFT
        # Настройка шрифта заголовка главы
        for run in chapter_heading.runs:
            run.font.name = 'Times New Roman'
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 0, 0)  # Черный цвет
        
        # Разбиваем текст на блоки по двойным переносам строк
        text_blocks = [block.strip() for block in chapter_text_cleaned.split('\n\n') if block.strip()]
        
        # Добавляем каждый блок как отдельный параграф
        for block in text_blocks:
            p = doc.add_paragraph(block)
            p.paragraph_format.line_spacing = 1.5
            p.paragraph_format.first_line_indent = Cm(1.25)
            p.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY  # Выравнивание по ширине
            p.paragraph_format.space_after = Pt(0)  # Убираем отступ после параграфа
            p.paragraph_format.space_before = Pt(0)  # Убираем отступ перед параграфом
            # Убеждаемся, что текст использует правильный шрифт
            for run in p.runs:
                run.font.name = 'Times New Roman'
                run.font.size = Pt(14)
        
        # Добавляем разрыв страницы после каждой главы (кроме последней)
        if i < len(chapters_text):
            doc.add_page_break()
    
    # === СПИСОК ИСТОЧНИКОВ ===
    
    # Определяем количество источников в зависимости от типа работы
    if work_type in ["Курсовая работа", "Дипломная работа"]:
        sources_count = 20
    else:
        sources_count = 12  # 10-15 источников, берем среднее значение
    
    # Получаем ключевые слова для поиска источников
    keywords = extract_keywords_from_theme(work_theme, science_name)
    
    # Получаем источники через Coze workflow
    sources = await fetch_sources_from_coze(keywords, sources_count)
    
    # Если источники получены, добавляем их в документ
    if sources:
        # Заголовок раздела
        sources_heading = doc.add_heading("СПИСОК ИСТОЧНИКОВ", level=1)
        sources_heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in sources_heading.runs:
            run.font.name = 'Times New Roman'
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 0, 0)  # Черный цвет
        
        # Добавляем каждый источник
        for i, source in enumerate(sources, 1):
            formatted_source = format_source_gost(source, i)
            source_p = doc.add_paragraph(formatted_source)
            source_p.paragraph_format.line_spacing = 1.5
            source_p.paragraph_format.first_line_indent = Cm(0)  # Без отступа для списка
            source_p.paragraph_format.left_indent = Cm(0)
            
            # Настройка шрифта
            for run in source_p.runs:
                run.font.name = 'Times New Roman'
                run.font.size = Pt(14)
        
        logging.info(f"Добавлено {len(sources)} источников в документ")
        
        # Уведомляем пользователя о завершении
        try:
            chat_id = get_chat_id(context)
            await context.bot.send_message(
                chat_id=chat_id,
                text=f"✅ Найдено и добавлено {len(sources)} источников"
            )
        except Exception as chat_error:
            logging.error(f"Ошибка отправки уведомления о завершении поиска: {chat_error}")
    else:
        logging.warning("Не удалось получить источники через Coze workflow")
        # Уведомляем пользователя
        try:
            chat_id = get_chat_id(context)
            await context.bot.send_message(
                chat_id=chat_id,
                text="⚠️ Не удалось автоматически найти источники. Пожалуйста, добавьте их самостоятельно."
            )
        except Exception as chat_error:
            logging.error(f"Ошибка отправки уведомления об ошибке поиска: {chat_error}")

    # Сохранение документа в память
    doc_io = io.BytesIO()
    doc.save(doc_io)
    doc_io.seek(0)  # Возвращаем указатель в начало
    
    logging.info("Документ создан в памяти")
    return doc_io

async def cancel(update: Update, context: CallbackContext) -> int:
    """Обработчик отмены заказа"""
    user_id = update.effective_user.id
    await log_user_action(user_id, "cancel_order")
    
    keyboard = [["/order"], ["/help", "/menu"]]
    reply_markup = ReplyKeyboardMarkup(keyboard, resize_keyboard=True)
    
    await update.message.reply_text(
        "❌ Заказ отменён.\n\n"
        "Вы можете начать новый заказ с помощью команды /order",
        reply_markup=reply_markup
    )
    return ConversationHandler.END

# Команда для заказа презентации
async def order_presentation(update: Update, context: CallbackContext) -> int:
    """Начинает процесс заказа презентации"""
    user_id = update.effective_user.id
    await log_user_action(user_id, "presentation_command")
    
    keyboard = [["/order"], ["/presentation"], ["/help", "/cancel"]]
    reply_markup = ReplyKeyboardMarkup(keyboard, resize_keyboard=True)
    
    await update.message.reply_text(
        "📊 *Заказ презентации PowerPoint* 📊\n\n"
        "Презентация будет создана автоматически на основе доклада.\n"
        "Включает:\n"
        "✅ Профессиональный дизайн\n"
        "✅ Изображения по теме\n"
        "✅ Графики и диаграммы\n"
        "✅ Структурированный контент\n\n"
        "💰 Стоимость: 400 рублей\n\n"
        "Для заказа используйте команду /order и выберите тип работы 'Доклад'.\n"
        "После генерации доклада вы получите также презентацию!",
        reply_markup=reply_markup,
        parse_mode='Markdown'
    )
    return ConversationHandler.END

# Основная функция
def main():
    application = ApplicationBuilder().token(TELEGRAM_BOT_TOKEN).build()

    conv_handler = ConversationHandler(
        entry_points=[CommandHandler("order", order)],
        states={
            WORK_TYPE: [MessageHandler(filters.TEXT & ~filters.COMMAND, work_type_handler)],
            SCIENCE_NAME: [MessageHandler(filters.TEXT & ~filters.COMMAND, science_name_handler)],
            PAGE_NUMBER: [MessageHandler(filters.TEXT & ~filters.COMMAND, page_number_handler)],
            WORK_THEME: [MessageHandler(filters.TEXT & ~filters.COMMAND, work_theme_handler)],
            CUSTOM_PLAN: [MessageHandler(filters.TEXT & ~filters.COMMAND, custom_plan_handler)],
            PREFERENCES: [MessageHandler(filters.TEXT & ~filters.COMMAND, preferences_handler)],
            PAYMENT: [
                CallbackQueryHandler(create_payment, pattern="^pay$"),
                CallbackQueryHandler(back_button_handler, pattern="^back$"),
                CommandHandler("cancel", cancel)
            ]
        },
        fallbacks=[CommandHandler("cancel", cancel)],
        allow_reentry=True
    )

    # Основные команды
    application.add_handler(CommandHandler("start", start))
    application.add_handler(CommandHandler("presentation", order_presentation))
    application.add_handler(MessageHandler(filters.Regex("^Продолжить$"), continue_handler))
    application.add_handler(CommandHandler("help", help_command))
    application.add_handler(CommandHandler("menu", menu))
    application.add_handler(conv_handler)
    
    # Админ-команды
    application.add_handler(CommandHandler("admin_stats", admin_stats))
    application.add_handler(CommandHandler("admin_users", admin_users))
    application.add_handler(CommandHandler("admin_orders", admin_orders))
    application.add_handler(CommandHandler("admin_system", admin_system))
    application.add_handler(CommandHandler("broadcast", admin_broadcast))
    application.add_handler(CommandHandler("admin_finance", admin_finance))
    application.add_handler(CommandHandler("admin_export", admin_export))
    application.add_handler(CommandHandler("admin_monitor", admin_monitor))

    # Вывод информации о режиме работы
    mode_indicator = "ТЕСТОВЫЙ РЕЖИМ (без реальных платежей)" if TESTING_MODE else "РАБОЧИЙ РЕЖИМ (с реальными платежами)"
    print("="*60)
    print(f"🚀 Бот запущен с улучшенной админ-панелью!")
    print(f"{'🧪' if TESTING_MODE else '💳'} {mode_indicator}")
    print("="*60)
    logging.info("Бот запущен с улучшенной админ-панелью")
    logging.info(f"{'='*60}")
    logging.info(f"{mode_indicator}")
    logging.info(f"{'='*60}")
    application.run_polling()

if __name__ == "__main__":
    main()